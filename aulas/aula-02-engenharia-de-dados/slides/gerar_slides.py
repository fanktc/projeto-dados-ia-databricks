#!/usr/bin/env python3
"""Slides da Noite 2 — setup inicial + os 6 passos.

Slide também é código. Os tokens de design abaixo foram extraídos do deck
"Lançamento Agosto - Jornada de dados - v1 - Aula 01.pptx", então este arquivo
e aquele desenham a mesma coisa.

O conteúdo vem de três fontes, não de invenção:
  · transcript da Noite 1 (design sprint, as 5 dinâmicas)
  · transcript do workshop "Databricks + IA: o workflow completo com Claude Code"
    (o setup: bundle init, uv, auth, gh, AI Dev Kit, MCP, guard rails)
  · a execução real dos 6 prompts, com os números medidos

Uso:  python gerar_slides.py        (precisa de python-pptx)
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── tokens do sistema de design da Noite 1 ────────────────────────────
FUNDO      = RGBColor(0x07, 0x0C, 0x16)
ACENTO     = RGBColor(0x3B, 0x9D, 0xF5)
ACENTO_CLR = RGBColor(0x7C, 0xC8, 0xFF)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA      = RGBColor(0x9B, 0xAA, 0xC0)
MARCA      = RGBColor(0x56, 0x6A, 0x85)
CARD       = RGBColor(0x0D, 0x15, 0x22)
CARD_BORDA = RGBColor(0x1F, 0x3A, 0x5A)
DESTAQUE   = RGBColor(0x14, 0x23, 0x37)
CODIGO     = RGBColor(0x7C, 0xC8, 0xFF)

LARGURA, ALTURA = Inches(40), Inches(22.5)

prs = Presentation()
prs.slide_width, prs.slide_height = LARGURA, ALTURA
BRANCA = prs.slide_layouts[6]


def nova():
    s = prs.slides.add_slide(BRANCA)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = FUNDO
    return s


def texto(s, txt, x, y, w, h, tam, cor, negrito=False, fonte="Arial",
          alinha=PP_ALIGN.LEFT, entrelinha=1.0):
    cx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, linha in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = entrelinha
        r = p.add_run()
        r.text = linha
        r.font.name = fonte
        r.font.size = Pt(tam)
        r.font.bold = negrito
        r.font.color.rgb = cor
    return cx


def caixa(s, x, y, w, h, destaque=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = DESTAQUE if destaque else CARD
    sh.line.color.rgb = ACENTO if destaque else CARD_BORDA
    sh.line.width = Pt(6.4 if destaque else 4.0)
    sh.shadow.inherit = False
    return sh


def assinatura(s):
    texto(s, "JORNADA DE DADOS", 1.8, 20.48, 12.0, 1.0, 32, MARCA, True)


# ── arquétipo 1: divisor ──────────────────────────────────────────────
def divisor(chapeu, titulo, linha):
    s = nova()
    barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.7), ALTURA)
    barra.fill.solid()
    barra.fill.fore_color.rgb = ACENTO
    barra.line.fill.background()
    barra.shadow.inherit = False
    texto(s, chapeu, 2.8, 6.2, 30.0, 2.0, 52, ACENTO, True)
    texto(s, titulo, 2.8, 8.0, 34.4, 6.0, 168, BRANCO, True, entrelinha=0.95)
    texto(s, linha, 2.8, 15.6, 32.0, 2.4, 60, CINZA)
    return s


# ── arquétipo 2: linhas em card ───────────────────────────────────────
# itens: (rotulo, corpo) ou (rotulo, corpo, destaque). De 3 a 6 linhas.
MEDIDAS = {3: (3.02, 0.42, 62, 50), 4: (2.55, 0.34, 54, 44),
           5: (2.05, 0.30, 46, 38), 6: (1.66, 0.27, 40, 33)}


def linhas(chapeu, titulo, itens, rodape=None, mono=False):
    s = nova()
    texto(s, chapeu, 1.8, 1.8, 32.0, 1.12, 44, ACENTO, True)
    texto(s, titulo, 1.8, 3.5, 36.4, 4.4, 128, BRANCO, True, entrelinha=0.95)

    altura, folga, tam_rot, tam_corpo = MEDIDAS[len(itens)]
    topo = 8.17
    for i, item in enumerate(itens):
        rotulo, corpo = item[0], item[1]
        destaque = len(item) > 2 and item[2]
        y = topo + i * (altura + folga)
        caixa(s, 1.77, y, 36.46, altura, destaque)
        texto(s, rotulo, 3.0, y + altura * 0.20, 10.4, altura, tam_rot,
              ACENTO if destaque else BRANCO, True,
              fonte="Courier New" if mono else "Arial")
        texto(s, corpo, 13.8, y + altura * 0.20, 23.2, altura, tam_corpo,
              CINZA, entrelinha=1.05)
    if rodape:
        texto(s, rodape, 1.8, 19.35, 36.0, 1.4, 46, ACENTO_CLR)
    assinatura(s)
    return s


# ── arquétipo 3: quatro números ───────────────────────────────────────
def numeros(chapeu, titulo, tiles):
    s = nova()
    texto(s, chapeu, 1.8, 1.8, 32.0, 1.12, 44, ACENTO, True)
    texto(s, titulo, 1.8, 3.5, 36.4, 4.4, 128, BRANCO, True, entrelinha=0.95)
    for i, (simbolo, valor, desc) in enumerate(tiles[:4]):
        x = 1.77 + i * 9.28
        caixa(s, x, 9.57, 8.62, 7.66)
        texto(s, simbolo, x + 0.91, 10.32, 6.8, 1.2, 44, ACENTO, True)
        texto(s, valor, x + 0.91, 11.46, 6.8, 2.3, 58, BRANCO, True, entrelinha=0.95)
        texto(s, desc, x + 0.91, 14.0, 6.8, 2.9, 42, CINZA, entrelinha=1.05)
    assinatura(s)
    return s



# ── arquétipo 4: o fluxo (o que vamos construir) ──────────────────────
def fluxo(chapeu, titulo, subtitulo, etapas, faixa, paineis):
    """etapas: (camada, titulo, corpo). paineis: (nome, [linhas], nota)."""
    s = nova()
    texto(s, chapeu, 1.8, 1.4, 32.0, 1.0, 40, ACENTO, True)
    texto(s, titulo, 1.8, 2.6, 36.4, 3.0, 96, BRANCO, True, entrelinha=0.95)
    texto(s, subtitulo, 1.8, 5.5, 36.4, 1.2, 46, CINZA)

    # a esteira de camadas
    largura, folga, topo, altura = 5.6, 0.55, 7.3, 6.0
    for i, (camada, tit, corpo) in enumerate(etapas):
        x = 1.8 + i * (largura + folga)
        ultimo = i == len(etapas) - 1
        caixa(s, x, topo, largura, altura, destaque=ultimo)
        texto(s, camada, x + 0.45, topo + 0.45, largura - 0.9, 0.9, 30, ACENTO, True)
        texto(s, tit, x + 0.45, topo + 1.5, largura - 0.9, 1.6, 34, BRANCO, True, entrelinha=0.95)
        texto(s, corpo, x + 0.45, topo + 3.15, largura - 0.9, 2.6, 24, CINZA, entrelinha=1.15)
        if not ultimo:
            texto(s, "→", x + largura, topo + altura / 2 - 0.55, folga, 1.1, 44,
                  ACENTO, True, alinha=PP_ALIGN.CENTER)

    # a faixa que atravessa tudo
    faixa_sh = caixa(s, 1.8, 13.75, 36.4, 1.15)
    texto(s, faixa, 1.8, 14.05, 36.4, 0.8, 30, ACENTO_CLR, True, alinha=PP_ALIGN.CENTER)

    # os três painéis de baixo
    pw, pfolga, py, ph = 11.6, 0.8, 15.4, 4.1
    for i, (nome, itens, nota) in enumerate(paineis[:3]):
        x = 1.8 + i * (pw + pfolga)
        caixa(s, x, py, pw, ph)
        texto(s, nome, x + 0.6, py + 0.35, pw - 1.2, 0.8, 28, ACENTO, True)
        texto(s, "\n".join("· " + i2 for i2 in itens), x + 0.6, py + 1.3, pw - 1.2, 2.0,
              26, CINZA, entrelinha=1.25)
        texto(s, nota, x + 0.6, py + ph - 0.95, pw - 1.2, 0.8, 24, MARCA, True)
    assinatura(s)
    return s


# ── arquétipo 5: quatro colunas de dor ────────────────────────────────
def colunas(chapeu, titulo, subtitulo, cols, saida_nome, saida_itens, saida_nota):
    """cols: (etiqueta, manchete, corpo)."""
    s = nova()
    texto(s, chapeu, 1.8, 1.4, 32.0, 1.0, 40, ACENTO, True)
    texto(s, titulo, 1.8, 2.6, 36.4, 3.4, 104, BRANCO, True, entrelinha=0.95)
    texto(s, subtitulo, 1.8, 6.3, 36.4, 1.2, 48, CINZA)

    cw, cfolga, cy, ch = 8.5, 0.8, 8.6, 7.4
    for i, (etiqueta, manchete, corpo) in enumerate(cols[:4]):
        x = 1.8 + i * (cw + cfolga)
        caixa(s, x, cy, cw, ch)
        texto(s, etiqueta, x + 0.7, cy + 0.6, cw - 1.4, 0.8, 30, ACENTO, True)
        texto(s, manchete, x + 0.7, cy + 1.7, cw - 1.4, 2.4, 50, BRANCO, True, entrelinha=0.95)
        texto(s, corpo, x + 0.7, cy + 4.5, cw - 1.4, 2.6, 32, CINZA, entrelinha=1.15)

    caixa(s, 1.8, 16.6, 36.4, 3.0, destaque=True)
    texto(s, saida_nome, 2.6, 17.0, 9.0, 0.9, 34, ACENTO, True)
    texto(s, "   ·   ".join(saida_itens), 2.6, 18.05, 26.0, 1.0, 44, BRANCO, True)
    texto(s, saida_nota, 29.0, 18.15, 8.4, 1.0, 30, ACENTO_CLR, True, alinha=PP_ALIGN.RIGHT)
    assinatura(s)
    return s


# ══════════════════════════════════════════════════════════════════════
#  PARTE 0 · A PONTE — DO DESIGN SPRINT PARA O CÓDIGO
# ══════════════════════════════════════════════════════════════════════

divisor("NOITE 2 · TERÇA 25/08",
        "Do design sprint\npara o código",
        "Ontem a gente descobriu o que construir. Hoje a gente constrói — em seis passos.")

# ── a dor que justifica a noite ───────────────────────────────────────
colunas("A DOR DE HOJE",
        "Databricks só no web\nworkspace não escala",
        "Se você desenvolve clicando no navegador, a conta chega depois — em produção.",
        [
          ("VERSIONAMENTO", "Notebook\nnão versiona",
           "Sem Git de verdade: sem PR, sem code review, sem histórico."),
          ("QUALIDADE", "Quebra em\nprodução",
           "Sem teste local nem CI — o erro só aparece rodando."),
          ("GOVERNANÇA", "Difícil\nauditar",
           "Quem mudou o quê? E o ambiente não é reproduzível."),
          ("IA", "IA presa\nno notebook",
           "Sem tools e sem o projeto todo. Vira copia-e-cola."),
        ],
        "A SAÍDA",
        ["Ambiente local", "Asset Bundles", "MCP", "AI Dev Kit"],
        "escala · audita · IA produtiva e segura")

# ── onde a gente vai chegar ───────────────────────────────────────────
fluxo("QUEM SABE FAZ AO VIVO",
      "O que vamos construir hoje",
      "Pipeline medallion + dashboard AI/BI + Genie na Rota do Perfume — tudo como código, dirigido por IA.",
      [
        ("FONTE", "ERP e CRM", "10 CSVs de origem\n14,7 MB · seed 42\ndado idêntico para todos"),
        ("RAW", "Arquivo no Volume", "bronze.raw governado\nconferência de chegada\n313.551 linhas"),
        ("BRONZE", "Ingestão bruta", "10 tabelas Delta\ntudo texto, sujeira intacta\n+ metadado de ingestão"),
        ("SILVER", "Limpo e conformado", "CNPJ, datas e dedup\nCONSTRAINT no Delta\ndevolução sinalizada"),
        ("GOLD", "Marts de negócio", "4 dimensões + fato\n3 marts por diretoria\n6 views de negócio"),
        ("CONSUMO", "Dashboard e Genie", "AI/BI como código\nGenie sobre dado limpo\nos dois dentro do bundle"),
      ],
      "UM JOB SERVERLESS · 12 TAREFAS · UM SCHEMA POR CAMADA · 11 TESTES QUE QUEBRAM O PIPELINE",
      [
        ("CLAUDE CODE",
         ["AI Dev Kit · 40 skills oficiais", "MCP · o workspace inteiro", "Asset Bundles (DABs)"],
         "a IA constrói · você valida"),
        ("GUARDRAILS",
         ["permissions deny no settings.json", "hook PreToolUse determinístico", "profile explícito, sempre"],
         "a IA não segura o botão"),
        ("ESTEIRA",
         ["bundle validate a cada passo", "11 testes que interrompem o job", "6 deploys, um por prompt"],
         "quebrou aqui, não quebra lá"),
      ])

linhas("RECAP DA NOITE 1", "As cinco dinâmicas\njá deram a resposta", [
    ("Como poderíamos…", "Todo problema virou pergunta, sem ninguém propor solução ainda."),
    ("Meta de 2 anos", "Escrita no passado, com número. Meta sem número é desejo."),
    ("O que pode dar errado", "Risco levantado ANTES de construir. Depois fica caro."),
    ("Quem faz o quê", "O caminho da pergunta até a ação — e o caminho de volta."),
    ("A menor coisa que responde", "O protótipo. Não é o produto final: é o suficiente para saber se presta.", True),
], rodape="\"Antes de escrever qualquer código, gaste um dia com o computador fechado, entrevistando o time.\"")

linhas("A PONTE", "O sprint não entrega\ncódigo. Entrega\nfeatures.", [
    ("O erro", "Jogar tudo num prompt só. Ele devolve alguma coisa — e você não sabe qual parte quebrou."),
    ("A saída certa", "De 5 a 6 features. Menos que 4 vira one-shot com outro nome; mais que 8 não cabe na noite."),
    ("Uma feature = um prompt", "Com começo, meio e fim. Você valida antes de seguir para a próxima.", True),
], rodape="\"Nunca sai perfeito de primeira. Faça uma vez e vá dando diretrizes.\" — por isso são seis, não um.")

divisor("PARTE 01", "Setup inicial",
        "Uma aula de Databricks sem abrir a tela do Databricks. Só o terminal.")

linhas("SETUP · O QUE PRECISA ESTAR INSTALADO", "Quatro pré-requisitos", [
    ("databricks", "A CLI do Databricks. É por ela que tudo acontece daqui para frente."),
    ("gh", "A CLI do GitHub. Leva o repositório local para o remoto em um comando."),
    ("uv", "O gerenciador de pacotes Python. Rápido, e resolve o ambiente sozinho."),
    ("claude", "O Claude Code. É o quarto ambiente da semana — o único que CONSTRÓI.", True),
], mono=True, rodape="Conta no Databricks Free Edition: sem cartão, sem pegadinha. Só serverless.")

linhas("SETUP · PASSO 1", "databricks\nbundle init", [
    ("Template", "Python. É o que traz src/, resources/, tests/ e o databricks.yml."),
    ("Jobs de exemplo", "Não. A gente vai escrever os nossos."),
    ("Serverless", "SIM — e não é preferência. O Free Edition só permite serverless.", True),
], mono=True, rodape="O Asset Bundle traz as boas práticas da engenharia de software para dentro da engenharia de dados.")

linhas("SETUP · PASSO 2", "O ambiente Python", [
    ("Python 3.12", "NÃO use 3.13. Bibliotecas do Databricks ainda quebram nela."),
    ("uv venv --seed", "O --seed já instala o pip por padrão, e evita dor de cabeça depois."),
    ("uv sync", "Instala o que o template do bundle já trouxe.", True),
], mono=True, rodape="É o erro número 1 de quem tenta em casa: instalou tudo no 3.13 e nada funciona.")

linhas("SETUP · PASSO 3", "databricks\nauth login", [
    ("O que ele faz", "Abre o navegador, você autentica, e pronto — o profile fica salvo."),
    ("O efeito colateral bom", "Ele preenche sozinho o host lá no databricks.yml do bundle."),
    ("Sempre com --profile", "Nunca deixe o profile implícito. Um dia você deploya em prod sem querer.", True),
], mono=True, rodape="Dev e prod isolados, testes unitários, esteira de CI/CD — o bundle usa Terraform por baixo.")

linhas("SETUP · PASSO 4", "git init\ngh repo create", [
    ("git init + commit", "O famoso first commit. Antes de a IA escrever a primeira linha."),
    ("gh repo create", "Leva o repositório local para o GitHub, sem abrir o navegador."),
    ("Por que agora", "Se a IA vai escrever código, você precisa de git diff e de git revert.", True),
], mono=True, rodape="O bundle já nasce com um README explicando como executar o projeto.")

linhas("SETUP · PASSO 5", "AI Dev Kit:\nskills e MCP", [
    ("Onde achar", "Pesquise \"AI Dev Kit Databricks\". É o repositório oficial de solutions."),
    ("40 skills", "Por grupo: engenharia de dados, BI, machine learning, app developer."),
    ("Escopo", "Instale no PROJETO, não global. Cada projeto tem o contexto dele."),
    ("Skill × MCP", "Skill é injeção de CONHECIMENTO. MCP é AÇÃO — é o que conecta no workspace.", True),
], rodape="Se o MCP não aparecer: reinicie o Claude Code. É quase sempre isso.")

linhas("SETUP · PASSO 6 — E O MAIS IMPORTANTE", "Guard rails\nANTES do MCP", [
    (".claude/settings.json", "Negue: bundle destroy, deploy --target prod, rm -rf, git push --force."),
    (".claude/hooks/", "Um script que bloqueia DROP, TRUNCATE e DELETE sem WHERE."),
    ("Por que hook", "Hook é DETERMINÍSTICO. Skill e MCP são probabilísticos: se ele bloqueia, bloqueia sempre.", True),
], rodape="\"Time de marketing plugou o MCP num lugar e fez besteira.\" O guard rail existe para isso não acontecer.")

# ══════════════════════════════════════════════════════════════════════
#  PARTE 2 · OS 6 PASSOS
# ══════════════════════════════════════════════════════════════════════

divisor("PARTE 02", "Os 6 passos",
        "Seis features, seis prompts, seis deploys. O bundle nasce vazio e termina completo.")

linhas("OS 6 PASSOS", "Uma feature por prompt", [
    ("1 · Raw", "Catálogo como código, Volume, e a conferência de chegada dos arquivos"),
    ("2 · Bronze", "Os 10 CSVs viram Delta. Nenhuma limpeza — a sujeira é preservada"),
    ("3 · Silver", "A limpeza, com o contrato de qualidade declarado na tabela"),
    ("4 · Gold", "Dimensões, fato, marts por diretoria e os 9 testes que quebram o job"),
    ("5 · Dashboard", "AI/BI como código, versionado dentro do bundle"),
    ("6 · Agentes de IA", "Metadado, views de negócio e o Genie sobre dado limpo", True),
], rodape="Cada passo termina em `databricks bundle deploy`. Deploy é rotina, não evento de fim de projeto.")

linhas("PASSO 1 · RAW", "O dado chega\nao Volume", [
    ("Raw não é bronze", "Raw é ARQUIVO. Bronze é TABELA. O CSV fica byte por byte — é a prova."),
    ("Volume, não DBFS", "Volume é objeto do Unity Catalog: tem dono, permissão e linhagem."),
    ("Conferência de chegada", "A tarefa mais chata do pipeline, e a que mais salva emprego.", True),
], rodape="Arquivo que não chega NÃO dá erro. Ele dá número menor — e o dashboard mente com cara de verdade.")

linhas("PASSO 2 · BRONZE", "Dez tabelas\nem um comando", [
    ("Ontem", "Vocês subiram tabela por tabela, clicando na interface."),
    ("Hoje", "Uma função e uma lista. Se o ERP mandar a 11ª tabela, é uma linha."),
    ("Tudo entra como texto", "Se inferir tipo, 15/10/2025 vira nulo e 309 CNPJs perdem o zero à esquerda.", True),
], rodape="A bronze preserva o problema para que ele possa ser RESOLVIDO — não escondido. E nunca se edita.")

linhas("PASSO 3 · SILVER", "A limpeza\ncom contrato", [
    ("A decisão da noite", "Quantidade negativa é devolução. Descartar infla o faturamento; sinalizar é o certo."),
    ("Deduplicar ≠ DISTINCT", "40 CNPJs com dois cadastros. O id é diferente, então DISTINCT não vê nada."),
    ("CONSTRAINT é contrato", "O Delta passa a RECUSAR a escrita que violar a regra. Ela vira da tabela, não do script.", True),
], rodape="A regra `valor_liquido >= 0` parecia óbvia e FALHOU em 135 pedidos. A regra errada era a nossa.")

linhas("PASSO 4 · GOLD", "Modelada para\nquem consome", [
    ("O contrato antes do SQL", "Granularidade, dimensões, métricas e filtros. Escritos numa frase."),
    ("Um fato, vários marts", "O que separa um mart do outro é a dimensão dominante, nunca a tabela base."),
    ("Teste que não quebra", "…não é teste, é relatório. Se falhar, o job PARA.", True),
], rodape="Melhor o dashboard ficar com o dado de ontem do que com o dado errado de hoje.")

linhas("PASSO 5 · DASHBOARD", "Também é\nartefato de código", [
    ("Clicado", "Não tem diff, não tem revisão, não tem rollback. Se alguém apaga, acabou."),
    ("Em JSON no bundle", "Sobe junto com o deploy. Desfazer é git revert."),
    ("Métrica declarada uma vez", "Nenhuma tela mostra receita diferente da outra. É o fim da reunião travada.", True),
], rodape="Ontem o dataset tinha CAST e dois try_to_date. Hoje é `receita`. É o que a silver comprou.")

linhas("PASSO 6 · AGENTES DE IA", "O mesmo dado,\noutra porta", [
    ("COMMENT é interface", "Não é documentação para humano. É o que o agente lê para DECIDIR a coluna."),
    ("View com nome de negócio", "Ninguém pergunta por fato_vendas. Perguntam por clientes em risco."),
    ("A regra que ele não adivinha", "Que dezembro é vale POR DESENHO do setor. Alguém precisa escrever.", True),
], rodape="Ontem o Genie foi plugado na bronze com o aviso de que podia errar. Hoje ele acerta — mesmo modelo.")

# ══════════════════════════════════════════════════════════════════════
#  FECHAMENTO
# ══════════════════════════════════════════════════════════════════════

numeros("O QUE OS 6 PASSOS ENTREGARAM", "Rodado do\ncatálogo vazio", [
    ("R$", "102.303.828,05", "a mesma receita da noite 1, nas três camadas"),
    ("191.080", "linhas", "no fato, grão de item de pedido"),
    ("12", "tarefas", "num job agendado, com 11 testes que quebram"),
    ("6", "deploys", "um por prompt. Nenhum deles falhou"),
])

divisor("O FECHAMENTO",
        "Engenharia de dados\nnão é o que a IA\nsubstitui",
        "É o que faz a IA funcionar. Amanhã: as três perguntas da diretoria viram modelo.")

# ══════════════════════════════════════════════════════════════════════
import os
saida = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                     "aula-02-engenharia-de-dados.pptx")
prs.save(saida)
print(f"{len(prs.slides._sldIdLst)} slides → {saida}")
