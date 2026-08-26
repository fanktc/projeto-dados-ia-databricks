#!/usr/bin/env python3
"""Slides da Noite 3 — "Quais 200?" · ciência de dados e agentes de IA.

Slide também é código. A noite inteira responde UMA pergunta:

    "Tenho 3.000 clientes. O time consegue ligar para 200 por semana.
     Quais 200?"

Três prompts ao vivo: features → modelo/MLflow → agente. Os divisores e os
slides de conceito são o que se fala ENQUANTO o Claude Code trabalha.

Formato 16:9 (13,333 x 7,5 pol). Paleta e regras de layout no fim do PRD.

Uso:  python gerar_slides.py        (precisa de python-pptx)
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── paleta ────────────────────────────────────────────────────────────
FUNDO      = RGBColor(0x07, 0x0C, 0x16)
CARD       = RGBColor(0x0D, 0x15, 0x22)
CARD_2     = RGBColor(0x14, 0x23, 0x37)
BORDA      = RGBColor(0x1F, 0x3A, 0x5A)
ACENTO     = RGBColor(0x3B, 0x9D, 0xF5)
ACENTO_CLR = RGBColor(0x7C, 0xC8, 0xFF)
ALERTA     = RGBColor(0xFF, 0x6B, 0x5A)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA      = RGBColor(0x9B, 0xAA, 0xC0)
MARCA      = RGBColor(0x56, 0x6A, 0x85)

SERIF = "Arial"
MONO  = "Consolas"

LARGURA, ALTURA = Inches(13.333), Inches(7.5)
MARGEM = 0.75
UTIL   = 13.333 - 2 * MARGEM      # 11,833

prs = Presentation()
prs.slide_width, prs.slide_height = LARGURA, ALTURA
BRANCA = prs.slide_layouts[6]


# ── primitivas ────────────────────────────────────────────────────────
def nova(cor=FUNDO):
    s = prs.slides.add_slide(BRANCA)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = cor
    return s


def texto(s, txt, x, y, w, h, tam, cor, negrito=False, fonte=SERIF,
          alinha=PP_ALIGN.LEFT, entrelinha=1.0, spc=None):
    cx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, linha in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = entrelinha
        r = p.add_run()
        r.text = linha
        r.font.name = fonte
        r.font.size = Pt(tam)
        r.font.bold = negrito
        r.font.color.rgb = cor
        if spc is not None:                     # letter-spacing, em centipontos
            r.font._rPr.set("spc", str(int(spc)))
    return cx


def cartao(s, x, y, w, h, tom=CARD, borda=BORDA, largura_borda=1.25):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = tom
    sh.line.color.rgb = borda
    sh.line.width = Pt(largura_borda)
    sh.shadow.inherit = False
    return sh


def barra(s, cor=ACENTO):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            Inches(0.22), ALTURA)
    sh.fill.solid()
    sh.fill.fore_color.rgb = cor
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def assinatura(s):
    texto(s, "JORNADA DE DADOS", MARGEM, 6.98, 6.0, 0.3, 9, MARCA, True, spc=120)


def nota(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def grade(y, h, n, folga=0.26):
    w = (UTIL - folga * (n - 1)) / n
    return [(MARGEM + i * (w + folga), y, w, h) for i in range(n)]


# ── arquétipos ────────────────────────────────────────────────────────
def cabecalho(kicker, titulo, sub=None, tam=36):
    """Kicker + título + subtítulo. Devolve (slide, y livre)."""
    s = nova()
    texto(s, kicker.upper(), MARGEM, 0.5, UTIL, 0.28, 11.5, ACENTO, True, spc=180)
    n = titulo.count("\n") + 1
    texto(s, titulo, MARGEM, 0.88, UTIL, 0.62 * n, tam, BRANCO, True,
          entrelinha=0.95, spc=-30)
    y = 0.88 + n * (tam / 72 * 1.06) + 0.24
    if sub:
        texto(s, sub, MARGEM, y, UTIL, 0.4, 14, CINZA, entrelinha=1.15)
        y += 0.34 * (sub.count("\n") + 1) + 0.22
    assinatura(s)
    return s, y


def divisor(kicker, titulo, sub):
    s = nova()
    barra(s)
    texto(s, kicker.upper(), 1.15, 2.55, 10.0, 0.32, 13, ACENTO, True, spc=200)
    n = titulo.count("\n") + 1
    texto(s, titulo, 1.15, 3.05, 11.3, 1.1 * n, 60, BRANCO, True,
          entrelinha=0.95, spc=-45)
    texto(s, sub, 1.15, 3.05 + n * 0.92 + 0.42, 10.6, 0.5, 18, CINZA)
    return s


def impacto(linhas):
    s = nova(ACENTO)
    texto(s, linhas, 1.3, 2.55, 10.73, 2.6, 34, BRANCO, True,
          alinha=PP_ALIGN.CENTER, entrelinha=1.35, spc=-20)
    return s


def bloco(s, y, txt, h=0.78, tom=CARD_2, borda=ACENTO, cor=BRANCO, tam=15,
          negrito=True, alinha=PP_ALIGN.LEFT):
    """Faixa de destaque, largura total."""
    cartao(s, MARGEM, y, UTIL, h, tom=tom, borda=borda, largura_borda=1.5)
    n = txt.count("\n") + 1
    texto(s, txt, MARGEM + 0.42, y + (h - n * tam / 72 * 1.28) / 2 + 0.02,
          UTIL - 0.84, h, tam, cor, negrito, alinha=alinha, entrelinha=1.28)


def codigo(s, y, txt, h=None, tam=12.5):
    n = txt.count("\n") + 1
    h = h or n * tam / 72 * 1.42 + 0.5
    cartao(s, MARGEM, y, UTIL, h, tom=CARD, borda=ACENTO, largura_borda=1.25)
    texto(s, txt, MARGEM + 0.4, y + 0.25, UTIL - 0.8, h, tam, ACENTO_CLR,
          fonte=MONO, entrelinha=1.42)
    return h


def mono(s, y, txt, h=None, tam=14, cor=ACENTO_CLR, recuo=None):
    """Diagrama em monoespaçada. recuo=x centraliza o BLOCO (não cada linha),
    que é o que mantém setas e conectores no lugar."""
    n = txt.count("\n") + 1
    h = h or n * tam / 72 * 1.5 + 0.55
    cartao(s, MARGEM, y, UTIL, h)
    if recuo is None:
        largura = max(len(l) for l in txt.split("\n")) * tam / 72 * 0.6
        recuo = max((UTIL - largura) / 2, 0.5)
    texto(s, txt, MARGEM + recuo, y + 0.28, UTIL - recuo, h, tam, cor, True,
          fonte=MONO, entrelinha=1.5)
    return h


def cartoes(s, y, h, itens, tam_rot=17, tam_corpo=12.5, mono_rot=False):
    """itens: (rotulo, corpo) ou (rotulo, corpo, 'alerta'|'destaque')."""
    for (x, yy, w, hh), item in zip(grade(y, h, len(itens)), itens):
        marca = item[2] if len(item) > 2 else None
        cor_b = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BORDA)
        tom   = CARD_2 if marca else CARD
        cartao(s, x, yy, w, hh, tom=tom, borda=cor_b,
               largura_borda=1.5 if marca else 1.25)
        cor_r = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BRANCO)
        nr = item[0].count("\n") + 1
        texto(s, item[0], x + 0.28, yy + 0.26, w - 0.56, 0.9, tam_rot, cor_r,
              True, fonte=MONO if mono_rot else SERIF, entrelinha=1.05)
        texto(s, item[1], x + 0.28, yy + 0.26 + nr * tam_rot / 72 * 1.22 + 0.16,
              w - 0.56, hh, tam_corpo, CINZA, entrelinha=1.28)


def faixas(s, y, itens, alt=0.62, folga=0.14, larg_esq=4.6, tam=13.5,
           cabecalho_=None):
    """Tabela de duas colunas em faixas. itens: (esq, dir) ou (esq, dir, marca)."""
    if cabecalho_:
        texto(s, cabecalho_[0], MARGEM + 0.34, y, larg_esq, 0.3, 11, ACENTO, True, spc=140)
        texto(s, cabecalho_[1], MARGEM + 0.34 + larg_esq, y, UTIL - larg_esq - 0.7,
              0.3, 11, ACENTO, True, spc=140)
        y += 0.42
    for item in itens:
        marca = item[2] if len(item) > 2 else None
        cor_b = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BORDA)
        tom   = CARD_2 if marca else CARD
        cartao(s, MARGEM, y, UTIL, alt, tom=tom, borda=cor_b,
               largura_borda=1.5 if marca else 1.25)
        cor_e = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BRANCO)
        texto(s, item[0], MARGEM + 0.34, y + (alt - tam / 72 * 1.3) / 2,
              larg_esq - 0.2, alt, tam, cor_e, True)
        texto(s, item[1], MARGEM + 0.34 + larg_esq, y + (alt - tam / 72 * 1.3) / 2,
              UTIL - larg_esq - 0.7, alt, tam, CINZA)
        y += alt + folga
    return y


def fluxo(s, y, etapas, h=0.95, quebrado=False, tam=13):
    """Caixas ligadas por seta. quebrado: a última vem em alerta, com ✗."""
    n = len(etapas)
    seta = 0.42
    w = (UTIL - seta * (n - 1)) / n
    for i, e in enumerate(etapas):
        x = MARGEM + i * (w + seta)
        ultimo = i == n - 1
        alvo = ultimo and quebrado
        destaque = ultimo and not quebrado
        cartao(s, x, y, w, h,
               tom=CARD_2 if (alvo or destaque) else CARD,
               borda=ALERTA if alvo else (ACENTO if destaque else BORDA),
               largura_borda=1.5 if (alvo or destaque) else 1.25)
        texto(s, e, x + 0.16, y + (h - tam / 72 * 1.3 * (e.count("\n") + 1)) / 2,
              w - 0.32, h, tam, ALERTA if alvo else (ACENTO_CLR if destaque else BRANCO),
              True, fonte=MONO, alinha=PP_ALIGN.CENTER, entrelinha=1.3)
        if not ultimo:
            texto(s, "✗" if (quebrado and i == n - 2) else "→",
                  x + w, y + h / 2 - 0.19, seta, 0.4, 17,
                  ALERTA if (quebrado and i == n - 2) else ACENTO, True,
                  alinha=PP_ALIGN.CENTER)



# ── arquétipo 6: tabela de dados ──────────────────────────────────────
def tabela(s, y, colunas, linhas, larguras=None, tam=13, alt=0.58,
           destaque=None, cor_destaque=None):
    """Uma tabela com dados reais. colunas: lista de rótulos. linhas: lista de
    listas. destaque: índice da linha a realçar."""
    n = len(colunas)
    larguras = larguras or [UTIL / n] * n
    xs, acum = [], MARGEM
    for w in larguras:
        xs.append(acum); acum += w

    for x, w, rot in zip(xs, larguras, colunas):
        texto(s, rot, x + 0.22, y, w - 0.3, 0.3, 11, ACENTO, True, spc=120)
    yy = y + 0.4
    for i, linha in enumerate(linhas):
        marcada = destaque is not None and i == destaque
        cor = cor_destaque or ACENTO
        cartao(s, MARGEM, yy, UTIL, alt,
               tom=CARD_2 if marcada else CARD,
               borda=cor if marcada else BORDA,
               largura_borda=1.5 if marcada else 1.25)
        for x, w, val in zip(xs, larguras, linha):
            negrito = marcada or x == xs[0]
            cor_txt = (cor if marcada and x != xs[0] else
                       (BRANCO if x == xs[0] else CINZA))
            texto(s, str(val), x + 0.22, yy + (alt - tam / 72 * 1.3) / 2,
                  w - 0.3, alt, tam, cor_txt, negrito,
                  fonte=MONO if x != xs[0] else SERIF)
        yy += alt + 0.1
    return yy


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 1 · A PERGUNTA DOS 200                            slides 1–7
# ══════════════════════════════════════════════════════════════════════

# ── 1 · capa ──────────────────────────────────────────────────────────
s = nova()
barra(s)
texto(s, "IMERSÃO AO VIVO · NOITE 3 DE 4 · 26/08 · 19H30",
      1.15, 2.25, 11.0, 0.32, 13, ACENTO, True, spc=200)
texto(s, "Quais 200?", 1.1, 2.72, 11.3, 1.9, 92, BRANCO, True, spc=-70)
texto(s, "Ciência de dados e agentes de IA",
      1.15, 4.68, 11.0, 0.5, 25, ACENTO_CLR, True, spc=-15)
texto(s, "Rota do Perfume · Databricks + MLflow + Claude Code",
      1.15, 6.55, 11.0, 0.35, 13, MARCA)

# ── 2 · a ligação de segunda-feira ────────────────────────────────────
s, y = cabecalho("Ontem à noite, depois da aula", "O diretor comercial ligou")
cartao(s, MARGEM, y, UTIL, 3.62, tom=CARD, borda=ACENTO, largura_borda=1.5)
texto(s, "“Luciano, a plataforma ficou boa. Dashboard bonito, o Genie responde tudo.",
      MARGEM + 0.55, y + 0.4, UTIL - 1.1, 0.5, 16, BRANCO, entrelinha=1.3)
texto(s, "Mas o meu time não vive de dashboard. Eu tenho 42 vendedores e, somando\ntodo mundo, dá umas 200 ligações boas por semana. No melhor mês.",
      MARGEM + 0.55, y + 1.08, UTIL - 1.1, 1.0, 16, CINZA, entrelinha=1.32)
texto(s, "Eu tenho 3.000 clientes na base.",
      MARGEM + 0.55, y + 2.06, UTIL - 1.1, 0.5, 16, CINZA)
texto(s, "Me diz quais 200.”",
      MARGEM + 0.55, y + 2.68, UTIL - 1.1, 0.6, 24, ACENTO_CLR, True, spc=-20)
bloco(s, 6.18, "Essa é a aula de hoje.", h=0.62, tam=16)
nota(s, "Deixa a citação respirar. Leia em voz alta, sem pressa, e pare depois "
        "de 'Me diz quais 200'. O silêncio aqui vale mais que qualquer slide.")

# ── 3 · a conta ───────────────────────────────────────────────────────
s, y = cabecalho("O tamanho do problema", "200 de 3.000")
for (x, yy, w, h), (valor, desc) in zip(
        grade(y + 0.12, 2.5, 3),
        [("3.000", "clientes na base"),
         ("200", "ligações por semana"),
         ("6,7%", "da base por semana")]):
    cartao(s, x, yy, w, h)
    texto(s, valor, x, yy + 0.62, w, 1.2, 58, BRANCO, True,
          alinha=PP_ALIGN.CENTER, spc=-45)
    texto(s, desc, x, yy + 1.72, w, 0.5, 14, CINZA, alinha=PP_ALIGN.CENTER)
bloco(s, 5.62, "Cada cliente é abordado, em média, uma vez a cada 15 semanas.\n"
                "Se o vendedor ligar para os 200 errados, a empresa perde quatro meses.",
      h=1.12, tam=16)

# ── 4 · por que o dashboard não responde ──────────────────────────────
s, y = cabecalho("O limite do que a gente já tem",
                 "Tudo que a gente construiu olha para trás")
faixas(s, y, [
    ("Quem comprou mais no ano",   "Quem vai comprar semana que vem"),
    ("Qual marca vendeu melhor",   "Para quem ligar amanhã de manhã"),
    ("Quanto faturamos em outubro", "Quem está prestes a sumir"),
    ("DESCRIÇÃO",                  "PRIORIZAÇÃO", "destaque"),
], alt=0.72, larg_esq=5.4, cabecalho_=("O QUE A GOLD RESPONDE", "O QUE O DIRETOR PRECISA"))
texto(s, "Dashboard mostra o passado com precisão. Ninguém liga para o passado.",
      MARGEM, 6.42, UTIL, 0.4, 15, ACENTO_CLR)

# ── 5 · as respostas erradas ──────────────────────────────────────────
s, y = cabecalho("O que a maioria das empresas faz",
                 "Três jeitos de escolher os 200 — todos ruins")
cartoes(s, y + 0.1, 2.45, [
    ("Os maiores\nclientes", "Eles já compram. Você gasta ligação com quem ia comprar de qualquer jeito."),
    ("Os que sumiram\nhá mais tempo", "Muitos já foram embora. Você gasta ligação em cliente perdido."),
    ("A intuição\ndo vendedor", "Funciona para os 20 que ele lembra. E os outros 2.980?"),
], tam_rot=19, tam_corpo=13.5)
bloco(s, 5.72, "Todos os três são defensáveis numa reunião.\n"
                "E todos os três desperdiçam ligação.", h=1.0, tam=16)

# ── 6 · a resposta certa ──────────────────────────────────────────────
s, y = cabecalho("O que vamos construir", "Ordenar por probabilidade de compra")
fluxo(s, y + 0.35, ["3.000\nclientes", "modelo de\npropensão", "score\nde 0 a 1",
                    "ordena", "pega os\n200"], h=1.28)
bloco(s, y + 2.25, "Não é classificar quem compra e quem não compra. É ORDENAR.\n"
                    "O vendedor não precisa de um sim ou não — precisa de uma fila.",
      h=1.15, tam=17)
nota(s, "Fala: 'E sim, se a pergunta fosse sobre desconto, o modelo certo seria "
        "outro, chamado uplift. Mas aqui a pergunta é priorização de esforço, e "
        "para isso propensão é exatamente a ferramenta certa.'")

# ── 7 · como saber se funcionou ───────────────────────────────────────
s, y = cabecalho("A métrica que importa", "Não é acurácia.\nÉ quantos dos 200 compraram.")
faixas(s, y, [
    ("Ligar aleatório",                  "20     ·  a taxa base da semana"),
    ("Ligar para quem sumiu há mais tempo", "0     ·  esses já foram embora", "alerta"),
    ("Ligar para os maiores",            "44"),
    ("Ligar para os 200 de maior score", "86     ·  4,25× o aleatório", "destaque"),
], alt=0.62, larg_esq=6.4,
   cabecalho_=("ESTRATÉGIA", "DOS 200 ABORDADOS, QUANTOS COMPRAM"))
bloco(s, 5.78, "Isso tem nome: LIFT. Quantas vezes melhor que o aleatório.\n"
                "É o número que você leva para a reunião — não o AUC.", h=1.0, tam=16)
nota(s, "Todos medidos no workspace com seed 42, corte 2026-08-01, janela de 7 "
        "dias, score out-of-fold. O zero da segunda linha é o melhor número da "
        "aula: a intuição mais comum do comercial acerta nenhum dos 200. E "
        "ordenar pelo atraso puro, sem modelo, acerta 1.")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 2 · O PLANO DA NOITE                            slides 8–9
# ══════════════════════════════════════════════════════════════════════

# ── 8 · três prompts, três deploys ────────────────────────────────────
s, y = cabecalho("Como a gente vai construir", "Três prompts, três deploys")
tabela(s, y,
       ["", "O QUE SAI DE CADA UM", "ONDE ISSO FICA"],
       [["PROMPT 1", "Uma tabela com uma linha por cliente", "gold.features_cliente"],
        ["PROMPT 2", "Uma nota de 0 a 1 para cada cliente", "gold.score_propensao"],
        ["PROMPT 3", "A lista dos 200, com nome e motivo", "gold.fila_semanal"]],
       larguras=[2.0, 5.9, 3.93], tam=14, alt=0.85, destaque=2)
bloco(s, y + 3.3, "Cada prompt termina em DEPLOY, e o de baixo é a resposta do diretor.\n"
                   "Não é protótipo que vira produção depois: já nasce no pipeline.",
      h=1.05, tam=16)
nota(s, "Deploy não é etapa de fim de projeto — é o que acontece toda vez que "
        "você termina alguma coisa. Por isso são três, e não um no final.")

# ── 9 · ML é camada, não projeto ──────────────────────────────────────
s, y = cabecalho("Onde esse código vai morar", "ML é camada, não projeto à parte",
                 "Nenhum repositório novo. As três tarefas entram no MESMO job de terça.")
mono(s, y, "raw → bronze → silver ×4 → dimensões → fato → marts → testes\n"
           "                                            ├→ métricas → auditoria\n"
           "                                            └→ ml_features → ml_modelo → ml_fila",
     tam=13, h=2.05)
faixas(s, y + 2.3, [
    ("12 → 15 tarefas",  "O job de ontem ganha três. Continua sendo um bundle run só."),
    ("Os mesmos testes", "Modelo ruim quebra o pipeline, igual a dado nulo."),
    ("A mesma auditoria", "Faltou COMMENT numa coluna de feature? O job para.", "destaque"),
], alt=0.64, larg_esq=4.4)
nota(s, "É assim que nasce o modelo que ninguém consegue colocar em produção: "
        "repositório novo, notebook solto, ambiente à parte. Aqui ML é mais uma "
        "camada do mesmo pipeline.")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 3 · A MATÉRIA-PRIMA: A GOLD DE ONTEM            slides 10–12
# ══════════════════════════════════════════════════════════════════════

# ── 10 · a gold, em uma query ─────────────────────────────────────────
s, y = cabecalho("O que a noite 2 deixou pronto",
                 "A matéria-prima já está na mesa")
codigo(s, y, "SELECT COUNT(*) AS linhas, ROUND(SUM(receita), 2) AS receita\n"
             "FROM   lakehouse_rotaperfume.gold.fato_vendas;", tam=14)
for (x, yy, w, h), (valor, desc) in zip(
        grade(y + 1.35, 1.5, 2, folga=0.4),
        [("191.080", "linhas no fato, no grão de item de pedido"),
         ("R$ 102.303.828,05", "de receita, com devolução já descontada")]):
    cartao(s, x, yy, w, h, tom=CARD_2, borda=ACENTO, largura_borda=1.5)
    texto(s, valor, x + 0.35, yy + 0.24, w - 0.7, 0.6, 26, ACENTO_CLR, True, spc=-20)
    texto(s, desc, x + 0.35, yy + 0.92, w - 0.7, 0.4, 12.5, CINZA)
faixas(s, y + 3.15, [
    ("4 dimensões + 1 fato", "cliente, produto, vendedor, calendário — e fato_vendas."),
    ("3 marts + 6 views",    "vendedor, produto, financeiro. E as views que o Genie já lê."),
    ("11 testes",            "que interrompem o job quando o número não fecha.", "destaque"),
], alt=0.6, larg_esq=5.0)
nota(s, "Abra isso no SQL Editor de verdade. São 30 segundos e é o que dá "
        "autoridade ao resto: nada do que vem hoje é slide, é tabela.")

# ── 11 · o que a limpeza tirou do caminho ─────────────────────────────
s, y = cabecalho("Por que essa tabela existe",
                 "O dado não chegou assim",
                 "As dez sujeiras da noite 2, nos números do dataset de vocês.")
faixas(s, y, [
    ("3.443 pedidos",  "12% das datas vieram em dd/mm/aaaa, o resto em ISO."),
    ("40 clientes",    "duplicados: id novo, mesmo CNPJ escrito de outro jeito."),
    ("1.111 + 223",    "CNPJ pontuado e CNPJ com espaço em volta. Três formatos ao todo."),
    ("957 pedidos",    "cancelados com valor_total = 0 — e o ITEM sem flag nenhuma.", "alerta"),
    ("2.327 itens",    "de devolução, gravados como quantidade negativa."),
], alt=0.58, larg_esq=3.6)
bloco(s, 6.02, "A bronze guardou tudo isso exatamente como veio.\n"
                "A silver é onde some — e é por isso que hoje dá para construir feature.",
      h=1.0, tam=15)

# ── 12 · o que a limpeza vale em dinheiro ─────────────────────────────
s, y = cabecalho("O que acontece se ninguém limpar",
                 "R$ 3,5 milhões de receita\nque nunca existiu")
codigo(s, y, "-- somando o item cru, como ele veio do ERP\n"
             "SELECT SUM(valor_bruto) FROM bronze.itens_pedido;   -- 105.890.448,42\n\n"
             "-- a gold, com os 957 pedidos cancelados fora\n"
             "SELECT SUM(receita)     FROM gold.fato_vendas;      -- 102.303.828,05", tam=12.5)
faixas(s, y + 2.35, [
    ("A diferença", "R$ 3.586.620,37 em itens de pedido que o cliente cancelou."),
    ("Por que passa", "O pedido diz valor_total = 0. O item não diz nada.", "alerta"),
    ("O que quebraria hoje", "'Quanto este cliente gasta' erraria em todo cliente que já cancelou.", "destaque"),
], alt=0.62, larg_esq=4.8)
nota(s, "Este é o slide que amarra a noite 2 na noite 3: a feature valor_total do "
        "prompt 1 sai daqui. Feature errada não dá erro — dá fila errada.")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 4 · AS PREMISSAS DE ML                          slides 13–15
# ══════════════════════════════════════════════════════════════════════

# ── 13 · o que é um modelo ────────────────────────────────────────────
s, y = cabecalho("Antes de qualquer código", "O que é um modelo, afinal")
faixas(s, y, [
    ("NÃO é",   "um programa com regras que alguém escreveu. Ninguém digita 'se atraso > 2 então quente'."),
    ("É",       "uma função que recebe colunas e devolve um número entre 0 e 1."),
    ("Aprende", "olhando milhares de clientes do passado em que a resposta JÁ é conhecida.", "destaque"),
    ("Devolve", "uma nota de propensão. Não uma certeza, não uma causa, não uma promessa."),
], alt=0.68, larg_esq=3.4)
bloco(s, y + 3.15, "Ele não descobre POR QUE o cliente compra.\n"
                    "Ele descobre quais combinações de colunas andaram junto com 'comprou'.",
      h=1.0, tam=16)

# ── 14 · o vocabulário ────────────────────────────────────────────────
s, y = cabecalho("Quatro palavras que vão aparecer a noite inteira",
                 "Feature, rótulo, treino e teste")
mono(s, y, "  FEATURES (X)                          RÓTULO (y)\n"
           "  o que se sabia até 31/07              o que aconteceu depois\n"
           "  ┌──────────────────────────┐          ┌──────────┐\n"
           "  │ recência  atraso  ticket │   ──→    │ comprou? │\n"
           "  └──────────────────────────┘          └──────────┘", tam=13, h=2.1)
faixas(s, y + 2.35, [
    ("Treino  ·  75%", "os clientes que o modelo VÊ. É olhando para eles que ele se ajusta."),
    ("Teste  ·  25%",  "os que ele NÃO vê. A única nota honesta que existe.", "destaque"),
], alt=0.66, larg_esq=4.0)
bloco(s, y + 3.9, "Nota tirada no treino é prova com o gabarito na mão.\n"
                    "Todo número que a gente levar para a tela sai do teste.", h=0.95, tam=15)

# ── 15 · as premissas ─────────────────────────────────────────────────
s, y = cabecalho("O que estamos assumindo — e onde isso quebra",
                 "As cinco premissas do projeto")
faixas(s, y, [
    ("O passado se repete", "setembro se parece o bastante com agosto. Quebra em mudança de política comercial."),
    ("O futuro não entrou", "nenhuma coluna sabe o que houve depois do corte. É o vazamento, e é o próximo slide."),
    ("Uma linha por cliente", "clientes não se influenciam. Quebra se dois forem a mesma rede."),
    ("Treino parecido com o score", "quem eu pontuo se parece com quem eu treinei. Quebra com base nova."),
    ("Ordenar, não explicar", "'atraso alto anda junto com compra' não é 'atraso alto CAUSA compra'.", "destaque"),
], alt=0.58, larg_esq=5.0)
texto(s, "Premissa não é detalhe teórico: é a lista do que conferir quando o "
         "modelo começar a errar.", MARGEM, 6.35, UTIL, 0.4, 15, ACENTO_CLR, True)
nota(s, "Se alguém perguntar 'e se a empresa mudar a política de desconto?': "
        "quebra a primeira premissa, o modelo precisa de retreino, e é exatamente "
        "por isso que ele mora num pipeline agendado e não num notebook.")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 5 · FEATURES  ·  prompt 1 rodando               slides 16–21
# ══════════════════════════════════════════════════════════════════════

divisor("Prompt 1 de 3 · primeiro passo", "O que descreve\num cliente",
        "A parte que mais vale dinheiro no projeto.")

# ── 17 · o que é feature engineering ──────────────────────────────────
s, y = cabecalho("A palavra que descreve 80% do trabalho",
                 "O que é feature engineering")
bloco(s, y, "Transformar o dado que existe em colunas que RESPONDEM à pergunta.\n"
             "A mesma informação, escrita de outro jeito, vale muito diferente.",
      h=1.05, tam=16)
cartoes(s, y + 1.25, 2.3, [
    ("AGREGAR",   "Muitas linhas viram uma.\n\n28.729 pedidos → um número\npor cliente: quantos, quanto,\nquando foi o último"),
    ("RELATIVIZAR", "Comparar o cliente com ele\nmesmo.\n\n'20 dias sem comprar' vira\n'2,9× o ciclo dele'", "destaque"),
    ("SINALIZAR", "Evento vira coluna.\n\n'apareceu na lista de\ncompradores do lançamento'\nvira 1 ou 0"),
], tam_rot=17, tam_corpo=12.5)
faixas(s, y + 3.8, [
    ("O que NÃO é", "escolher algoritmo, ajustar hiperparâmetro, procurar biblioteca nova."),
], alt=0.62, larg_esq=3.4)
nota(s, "É a parte que não dá para terceirizar: depende de saber como a empresa "
        "vende. Duas pessoas com o mesmo dado e o mesmo algoritmo chegam a "
        "modelos muito diferentes por causa deste slide.")

# ── 18 · de onde vêm as features ──────────────────────────────────────
s, y = cabecalho("Tudo sai da gold de ontem", "Nenhuma fonte nova")
mono(s, y + 0.3,
     "gold.fato_vendas       ─┐\n"
     "silver.oportunidades   ─┼→   gold.features_cliente\n"
     "silver.visitas         ─┘", tam=19, h=2.5)
texto(s, "Uma linha por cliente. Vinte colunas descrevendo o comportamento dele.\n"
         "É isso que o modelo vai ler.",
      MARGEM, 5.9, UTIL, 0.9, 16, ACENTO_CLR, entrelinha=1.3)

# ── 10 · RFM ──────────────────────────────────────────────────────────
s, y = cabecalho("O ponto de partida", "Recência, frequência, valor")
for (x, yy, w, h), (sig, perg, col) in zip(
        grade(y + 0.2, 2.6, 3),
        [("R", "Faz quanto tempo\nque ele comprou?", "recencia_dias"),
         ("F", "Com que frequência\nele compra?", "frequencia_pedidos"),
         ("V", "Quanto ele gasta?", "valor_total")]):
    cartao(s, x, yy, w, h)
    texto(s, sig, x + 0.32, yy + 0.3, w - 0.64, 0.8, 44, ACENTO, True, spc=-30)
    texto(s, perg, x + 0.32, yy + 1.12, w - 0.64, 1.0, 15, BRANCO, entrelinha=1.25)
    texto(s, col, x + 0.32, yy + 2.05, w - 0.64, 0.4, 13.5, ACENTO_CLR, True, fonte=MONO)
bloco(s, 5.98, "Três colunas. Se você só tiver isso, já é melhor que a intuição.",
      h=0.72, tam=16)

# ── 11 · RFM sozinho mente ────────────────────────────────────────────
s, y = cabecalho("O problema do básico", "Dois clientes, a mesma recência")
for (x, yy, w, h), (nome, linhas_, veredito, acao, alerta) in zip(
        grade(y + 0.1, 2.42, 2, folga=0.4),
        [("Perfumaria Aurora", "Comprou há 20 dias\nCompra a cada 7 dias",
          "Quase 3x atrasado", "Liga hoje", True),
         ("Boutique Essenza", "Comprou há 20 dias\nCompra a cada 90 dias",
          "Está no ritmo", "Não incomoda", False)]):
    cor = ALERTA if alerta else ACENTO
    cartao(s, x, yy, w, h, tom=CARD_2, borda=cor, largura_borda=1.5)
    texto(s, nome, x + 0.4, yy + 0.28, w - 0.8, 0.5, 21, BRANCO, True, spc=-20)
    texto(s, linhas_, x + 0.4, yy + 0.88, w - 0.8, 0.9, 14.5, CINZA, entrelinha=1.35)
    texto(s, veredito, x + 0.4, yy + 1.62, w - 0.8, 0.45, 19, cor, True, spc=-15)
    texto(s, acao, x + 0.4, yy + 2.02, w - 0.8, 0.35, 14, CINZA)
bloco(s, 5.68, "Mesma recência. Situações opostas.\n"
                "Ordenar por recência coloca os dois na mesma posição da fila.",
      h=1.05, tom=CARD_2, borda=ALERTA, cor=BRANCO, tam=16)

# ── NOVO · uma linha por cliente, com dado de verdade ─────────────────
s, y = cabecalho("O que a tabela de features tem dentro",
                 "Uma linha por cliente. Só isso.",
                 "Três clientes reais da base, com as colunas que o modelo vai ler.")
tabela(s, y,
       ["CLIENTE", "COMPROU HÁ", "COMPRA A CADA", "ATRASO", "PEDIDOS", "TOTAL"],
       [["Casa de Fragrâncias Sublime", "10 dias", "50 dias", "0,2", "15", "R$ 38 mil"],
        ["Bella Diva", "29 dias", "30 dias", "0,95", "11", "R$ 52 mil"],
        ["Comercial Lumiar", "545 dias", "43 dias", "10,0", "9", "R$ 61 mil"]],
       larguras=[3.9, 1.7, 2.0, 1.3, 1.4, 1.53], destaque=2, cor_destaque=ALERTA)
bloco(s, y + 2.5, "São 20 colunas assim. Nenhuma delas é opinião —\n"
                   "todas saem de somar e dividir o que já aconteceu.", h=1.0, tam=16)
nota(s, "O terceiro cliente está há 545 dias sem comprar, com ciclo de 43. Dez "
        "vezes atrasado. Pergunte para a sala se vale ligar para ele — a "
        "resposta intuitiva é sim, e o slide seguinte mostra que é não.")

# ── 12 · a feature que muda o jogo ────────────────────────────────────
s, y = cabecalho("Conhecimento de negócio virando coluna", "Atraso relativo")
cartao(s, MARGEM, y, UTIL, 1.0, tom=CARD, borda=ACENTO, largura_borda=1.5)
texto(s, "atraso_relativo  =  recencia_dias / intervalo_medio_dias",
      MARGEM, y + 0.3, UTIL, 0.5, 21, ACENTO_CLR, True, fonte=MONO,
      alinha=PP_ALIGN.CENTER)
cartoes(s, y + 1.2, 1.72, [
    ("~1,0",  "Está no ritmo dele\nROTINA"),
    ("2,0",   "Duas vezes atrasado\nATENÇÃO"),
    ("> 3,0", "Provavelmente foi pro concorrente\nURGENTE", "alerta"),
], tam_rot=26, tam_corpo=13.5, mono_rot=True)
texto(s, "Não vem de biblioteca nenhuma. Vem de entender como a empresa vende.\n"
         "É essa feature que vai ordenar a fila dos 200.",
      MARGEM, 6.02, UTIL, 0.9, 16, ACENTO_CLR, True, entrelinha=1.3)

# ── 13 · as 20 features ───────────────────────────────────────────────
s, y = cabecalho("O que o Claude Code vai construir", "Quatro grupos, uma tabela")
cartoes(s, y, 1.72, [
    ("RFM",   "recência, frequência, valor total, ticket médio, margem"),
    ("Ritmo", "intervalo médio, desvio do intervalo, atraso relativo", "destaque"),
    ("CRM",   "oportunidades, taxa de ganho, visitas, conversão de visita"),
    ("Mix",   "SKUs, categorias, concentração em marca, comprou lançamento"),
], tam_rot=17, tam_corpo=12)
codigo(s, y + 1.98,
       "SELECT razao_social, recencia_dias, intervalo_medio_dias,\n"
       "       ROUND(atraso_relativo, 1) AS atraso\n"
       "FROM   lakehouse_rotaperfume.gold.features_cliente\n"
       "ORDER  BY atraso_relativo DESC LIMIT 10;")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 6 · O MODELO, O AUC E O VAZAMENTO            slides 22–28
# ══════════════════════════════════════════════════════════════════════

divisor("Prompt 2 de 3 · segundo passo", "O modelo",
        "E o erro que mata projeto em produção.")

# ── NOVO · treino e teste, com os números reais ───────────────────────
s, y = cabecalho("Como se sabe se o modelo é bom",
                 "Esconda parte dos clientes.\nDepois confira.")
mono(s, y, "2.815 clientes que já compraram\n"
           "  │\n"
           "  ├──  2.111  TREINO   o modelo vê. É olhando para eles que aprende\n"
           "  │\n"
           "  └──    704  TESTE    o modelo NÃO vê. Guardados no cofre", tam=13, h=2.3)
faixas(s, y + 2.55, [
    ("Por que esconder", "Nota tirada no treino é prova com o gabarito na mão."),
    ("O que se mede", "Dos 704 escondidos, o modelo acertou quem ia comprar?", "destaque"),
], alt=0.66, larg_esq=4.0)
bloco(s, y + 4.15, "Todo número que a gente levar para a tela hoje sai dos 704.",
      h=0.7, tam=16)

# ── 15 · o problema, escrito direito ──────────────────────────────────
s, y = cabecalho("Especificação", "Antes de codar, escreva em uma frase")
cartao(s, MARGEM, y, UTIL, 1.42, tom=CARD_2, borda=ACENTO, largura_borda=1.5)
texto(s, "Dado o comportamento de um cliente até hoje, qual a probabilidade\n"
         "de ele fazer um pedido nos próximos 7 dias?",
      MARGEM + 0.5, y + 0.36, UTIL - 1.0, 1.0, 21, BRANCO, True,
      entrelinha=1.3, spc=-18)
cartoes(s, y + 1.72, 1.55, [
    ("Unidade", "Um cliente"),
    ("Janela",  "7 dias — a semana da ligação"),
    ("Rótulo",  "Comprou (1) ou não (0)"),
    ("Uso",     "Ordenar a fila dos 200", "destaque"),
], tam_rot=16, tam_corpo=13.5)

# ── 18 · o algoritmo, que é a parte fácil ─────────────────────────────
s, y = cabecalho("A parte que todo mundo acha que é o trabalho",
                 "O algoritmo tem três linhas")
codigo(s, y, "modelo = HistGradientBoostingClassifier(random_state=42)\n"
             "modelo.fit(X_treino, y_treino)\n"
             "score  = modelo.predict_proba(X_cliente)[:, 1]", tam=15)
faixas(s, y + 1.5, [
    ("Por que árvore, e não regressão",
     "Ela acha sozinha que atraso perto de 1,5 é ouro e acima de 3 é cliente perdido."),
    ("predict_proba, nunca predict",
     "predict devolve 0 ou 1. A fila precisa de nota, não de sim ou não.", "destaque"),
    ("random_state=42",
     "Sem semente fixa, dois treinos dão listas diferentes. E aí ninguém confia."),
], alt=0.7, larg_esq=5.6)
bloco(s, y + 3.9, "Estas três linhas são iguais em qualquer empresa do mundo.\n"
                   "O que muda é o X — e o X foi o prompt 1.", h=1.0, tam=16)
nota(s, "É aqui que a sala espera a aula ficar difícil, e é o slide mais curto "
        "da noite. Diga isso em voz alta.")

# ── NOVO · quem comprou ontem vale zero ───────────────────────────────
s, y = cabecalho("Dois clientes reais, dois extremos",
                 "Quem comprou ontem\nnão compra hoje")
tabela(s, y,
       ["CLIENTE", "COMPROU HÁ", "COMPRA A CADA", "NOTA DO MODELO"],
       [["Bella Diva", "29 dias", "30 dias", "0,97"],
        ["Comercial Sublime", "1 dia", "51 dias", "0,00"]],
       larguras=[4.2, 2.4, 2.6, 2.63], tam=15, alt=0.78)
bloco(s, y + 2.2,
      "O que comprou ONTEM vale zero. O que está há um mês, no ritmo dele,\n"
      "vale 0,97. É por isso que \"ligue para quem sumiu\" não funciona.",
      h=1.1, tam=17)
texto(s, "Distribuição funciona por ciclo de reposição: o varejista acabou de "
         "receber a mercadoria — ele não compra de novo agora.",
      MARGEM, y + 3.5, UTIL, 0.9, 15, ACENTO_CLR, entrelinha=1.3)
nota(s, "Este é o slide para parar e deixar a sala olhar. Os dois números saíram "
        "do modelo treinado, não são exemplo inventado.")

# ── 26 · o que é AUC, e a régua da sala ───────────────────────────────
s, y = cabecalho("A palavra que vai aparecer a noite inteira", "O que é AUC")
cartao(s, MARGEM, y, UTIL, 1.32, tom=CARD_2, borda=ACENTO, largura_borda=1.5)
texto(s, "Sorteie dois clientes: um que comprou e um que não.\n"
         "Com que frequência o modelo dá a nota maior para quem comprou?",
      MARGEM + 0.5, y + 0.32, UTIL - 1.0, 1.0, 18, BRANCO, True, entrelinha=1.3)
for (x, yy, w, h), (valor, quem, desc, marca) in zip(
        grade(y + 1.6, 2.4, 4),
        [("0,35", "quem sumiu", "a fila está INVERTIDA:\nseguir o contrário\nseria melhor", "alerta"),
         ("0,50", "a moeda", "acerta a ordem\nmetade das vezes.\nnão sabe nada", None),
         ("0,64", "os maiores", "acerta 64 de cada\n100 pares. melhor\nque nada", None),
         ("0,88", "o modelo", "quem comprou vem\nna frente em 88\nde cada 100 pares", "destaque")]):
    cor = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BORDA)
    cartao(s, x, yy, w, h, tom=CARD_2 if marca else CARD, borda=cor,
           largura_borda=1.5 if marca else 1.25)
    texto(s, valor, x + 0.3, yy + 0.26, w - 0.6, 0.7, 34,
          ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BRANCO),
          True, spc=-25)
    texto(s, quem, x + 0.3, yy + 0.92, w - 0.6, 0.35, 14, BRANCO, True)
    texto(s, desc, x + 0.3, yy + 1.32, w - 0.6, 1.0, 12, CINZA, entrelinha=1.3)
bloco(s, y + 4.25, "AUC compara FILAS. Lift conta LIGAÇÕES.\n"
                    "A primeira é para você conferir; a segunda é para a reunião.",
      h=1.0, tam=16)
nota(s, "Os quatro números são as respostas da sala, medidas. O 0,37 é o momento "
        "da noite: não é 'um pouco ruim', é invertido — a mesma coisa que o "
        "'0 de 200' do slide 7 diz em linguagem comercial. Se alguém perguntar "
        "por que não usar acurácia: com 10% de compradores, um modelo que diz "
        "'ninguém compra' acerta 90% e é inútil.")

# ── 27 · vazamento de dado ────────────────────────────────────────────
s, y = cabecalho("O erro que mais mata modelo na vida real",
                 "Prever a semana que vem\ncom informação da semana que vem")
mono(s, y, "|---- features até 31/07 ----|CORTE|---- comprou até 07/08? ----|\n"
           "                              01/08/2026", tam=13, h=1.28)
cartoes(s, y + 1.5, 1.62, [
    ("Errado", "Features com todo o histórico\nAUC 0,98 — e quebra em produção", "alerta"),
    ("Certo",  "Features cortadas em 31/07\nAUC menor — e funciona", "destaque"),
], tam_rot=18, tam_corpo=13.5)
bloco(s, y + 3.4, "AUC alto demais é sinal de vazamento, não de competência.",
      h=0.72, tom=CARD_2, borda=ALERTA, tam=16)

# ── 28 · o teste que quase ninguém escreve ────────────────────────────
s, y = cabecalho("Por que o modelo entra nos mesmos testes de ontem",
                 "Um dado errado quebra.\nUm modelo ruim funciona.")
faixas(s, y, [
    ("Dado nulo",  "Explode, fica vermelho, alguém é avisado no mesmo dia."),
    ("Modelo ruim", "Devolve nota para todo mundo, na faixa certa, sem erro nenhum."),
    ("O resultado", "Pipeline verde, dashboard atualizado, e a lista errada por seis meses.", "alerta"),
], alt=0.68, larg_esq=4.4)
codigo(s, y + 2.5,
       "assert auc > melhor_baseline + 0.05, \"o modelo não ganha da regra simples\"\n"
       "assert auc < 0.99,                   \"bom demais é vazamento\"\n"
       "assert lift_top200 >= 2.5,           \"a fila não se paga\"", tam=12.5)
texto(s, "O teste que quase ninguém escreve é o primeiro: o modelo ganha do que "
         "a gente já fazia de graça?", MARGEM, 6.52, UTIL, 0.4, 14, ACENTO_CLR, True)
nota(s, "Quebre um destes de propósito ao vivo. Teste que ninguém viu falhar não "
        "convence ninguém.")

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 7 · MLFLOW  ·  prompt 2 rodando                slides 29–32
# ══════════════════════════════════════════════════════════════════════

divisor("Ainda no prompt 2 · terceiro passo", "MLflow",
        "Por que não basta salvar um .pkl numa pasta.")

# ── 31 · o que é MLflow ───────────────────────────────────────────────
s, y = cabecalho("Antes do porquê, o quê", "O que é MLflow")
bloco(s, y, "Uma biblioteca aberta — criada pela Databricks — que registra tudo o que\n"
             "acontece num treino, e um servidor que guarda isso para sempre.",
      h=1.05, tam=16)
cartoes(s, y + 1.25, 2.15, [
    ("Tracking",  "O diário.\nCada treino vira um run,\ncom o que entrou e o\nque saiu."),
    ("Models",    "O formato do pacote.\nO modelo + o ambiente +\no schema de entrada,\nnum diretório só."),
    ("Registry",  "O versionamento.\nv1, v2, v3 e os apelidos\n@prod e @challenger.", "destaque"),
    ("Evaluation", "A comparação.\nDois modelos lado a lado,\nna mesma tela, com as\nmesmas métricas."),
], tam_rot=17, tam_corpo=12)
faixas(s, y + 3.62, [
    ("No Databricks", "já vem instalado, e o Registry É o Unity Catalog — o mesmo do dado.", "destaque"),
], alt=0.62, larg_esq=3.4)

# ── 32 · a pergunta de daqui a seis meses ─────────────────────────────
s, y = cabecalho("Quando você não estiver mais no projeto",
                 "“Esse modelo ainda está bom?”")
faixas(s, y, [
    ("Qual versão está em produção?",          "Ninguém sabe"),
    ("Com que dado foi treinada?",             "Perdido"),
    ("Está melhor ou pior que a anterior?",    "Impossível comparar"),
    ("Quem treinou e quando?",                 "O cara saiu da empresa"),
], alt=0.66, larg_esq=6.2, cabecalho_=("PERGUNTA", "SEM MLFLOW"))
bloco(s, 6.0, "O MLflow não faz a empresa vender mais. O MODELO faz.\n"
               "O MLflow garante que ele continue existindo depois que você sair.",
      h=1.02, tam=16)

# ── 19 · o que o MLflow resolve ───────────────────────────────────────
s, y = cabecalho("Três coisas", "O que o MLflow resolve")
cartoes(s, y, 1.42, [
    ("Tracking",       "Registra cada treino: parâmetros, métricas, artefatos"),
    ("Model Registry", "Versiona e controla o estágio: staging, produção"),
    ("Comparação",     "Dois modelos lado a lado, na mesma tela"),
], tam_rot=17, tam_corpo=12.5)
codigo(s, y + 1.66,
       "with mlflow.start_run(run_name=\"propensao_v1\"):\n"
       "    mlflow.log_params(params)\n"
       "    mlflow.log_metric(\"auc\", auc)\n"
       "    mlflow.log_metric(\"lift_top200\", lift)\n"
       "    mlflow.sklearn.log_model(modelo, \"modelo\",\n"
       "        registered_model_name=\"propensao_compra\")", tam=12)
texto(s, "Repare no lift_top200. Essa é a métrica que responde a pergunta do "
         "diretor — e é ela que a gente versiona.",
      MARGEM, 6.62, UTIL, 0.4, 13.5, ACENTO_CLR)

# ── 35 · por que usamos MLflow aqui ───────────────────────────────────
s, y = cabecalho("Por que isso importa NESTE projeto",
                 "Três motivos concretos,\nnão boas práticas")
faixas(s, y, [
    ("O job retreina sozinho",
     "Toda semana nasce um modelo novo. Sem registry, ninguém sabe qual gerou a fila de hoje."),
    ("A fila guarda a versão",
     "score_propensao grava o número da versão. Dá para auditar a lista de três semanas atrás.", "destaque"),
    ("Rollback é uma linha",
     "set_registered_model_alias(modelo, 'prod', 1). Nenhum deploy, nenhum código alterado."),
], alt=0.78, larg_esq=4.8)
codigo(s, y + 2.9, "# quem consome nunca precisa saber o número da versão\n"
                    "mlflow.sklearn.load_model(\"models:/...propensao_compra@prod\")", tam=13)
texto(s, "Sem isso, \"o modelo está pior esse mês\" é uma conversa sem prova de "
         "nenhum dos dois lados.", MARGEM, 6.6, UTIL, 0.4, 13.5, ACENTO_CLR)

# ══════════════════════════════════════════════════════════════════════
#  BLOCO 8 · A FILA  ·  prompt 3 rodando                slides 33–37
# ══════════════════════════════════════════════════════════════════════

divisor("Prompt 3 de 3 · quarto passo", "Os 200", "Do score à ligação.")

# ── 21 · o gap que mata projetos ──────────────────────────────────────
s, y = cabecalho("Onde os projetos de ML morrem",
                 "O último metro é o mais difícil")
fluxo(s, y + 0.15, ["Dado", "Modelo", "Score na\ntabela", "Alguém\nligar"],
      h=1.05, quebrado=True)
cartao(s, MARGEM, y + 1.55, UTIL, 1.75, tom=CARD, borda=ACENTO, largura_borda=1.5)
texto(s, "“Já vi empresa com modelo excelente rodando há dois anos, AUC de 0,89,\n"
         "e o vendedor continuando a ligar pela intuição. Porque o score estava\n"
         "numa tabela que ele nunca abriu.”",
      MARGEM + 0.5, y + 1.85, UTIL - 1.0, 1.3, 16, BRANCO, entrelinha=1.35)
texto(s, "O último metro é o mais difícil. E é o que o agente resolve.",
      MARGEM, y + 3.55, UTIL, 0.4, 16, ACENTO_CLR, True)

# ── NOVO · a fila que saiu do banco ───────────────────────────────────
s, y = cabecalho("O que o pipeline entregou",
                 "A lista de segunda-feira",
                 "Três linhas da fila_semanal, exatamente como saíram da tabela.")
tabela(s, y,
       ["CLIENTE", "NOTA", "POR QUE ESTÁ NA LISTA"],
       [["Essência Nova Era", "0,90", "Comprou lançamento recente. Alta chance de repetir."],
        ["Perfumaria Aurora", "0,88", "4 pedidos nos últimos 90 dias. Ciclo curto — não esfrie."],
        ["Bella Vita", "0,87", "Mais da metade das visitas viram pedido (67%). Vale a ida."]],
       larguras=[3.3, 1.2, 7.33], tam=12.5, alt=0.66)
bloco(s, y + 2.7, "O vendedor não recebe um número. Recebe um nome, uma ordem\n"
                   "e uma frase que ele entende antes de discar.", h=1.0, tam=16)

# ── NOVO · onde a lista aparece ───────────────────────────────────────
s, y = cabecalho("A pergunta que todo mundo faz",
                 "E onde eu vejo esses 200?")
cartoes(s, y, 2.5, [
    ("1 · No dashboard",
     "A aba 'Fila da semana'.\nO vendedor escolhe o\nnome dele no filtro e vê\na lista em ordem.\n\nÉ a tela dele.", "destaque"),
    ("2 · Numa query",
     "SELECT * FROM\nfila_semanal\nWHERE vendedor = '...'\n\nÉ uma tabela como\nqualquer outra."),
    ("3 · Perguntando",
     "'Quem eu ligo essa\nsemana?'\n\nO Genie responde em\nportuguês — e mostra\no SQL que gerou."),
], tam_rot=18, tam_corpo=12.5)
bloco(s, y + 2.75, "As três leem a MESMA tabela. Nenhuma delas tem número inventado:\n"
                    "o que muda é só a porta de entrada.", h=1.05, tam=16)
nota(s, "Abra o dashboard primeiro — é o que responde a pergunta do diretor. "
        "A query é o plano B se algo não carregar ao vivo, e está pronta no "
        "passo-a-passo. O Genie fica por último, porque é o mais impressionante "
        "e o menos previsível.")

# ── 22 · o agente ─────────────────────────────────────────────────────
s, y = cabecalho("A última camada", "Ele não inventa. Ele consulta.")
faixas(s, y, [
    ("priorizar_carteira",      "Cruza carteira + score + features e ordena"),
    ("contexto_cliente",        "Histórico, ticket médio, marcas preferidas"),
    ("sugerir_produtos",        "O que ele compra e parou de comprar"),
    ("checar_disponibilidade",  "Estoque real, com flag de ruptura"),
], alt=0.68, larg_esq=5.0, cabecalho_=("FERRAMENTA", "O QUE FAZ"))
bloco(s, 6.02, "Regra do sistema: use sempre as ferramentas, nunca invente número.\n"
                "Agente sem dado organizado por trás é chute com sotaque.",
      h=1.0, tam=16)
# rótulos das ferramentas em monoespaçada
for i, sh in enumerate(s.shapes):
    pass

# ── 23 · a resposta para o diretor ────────────────────────────────────
s, y = cabecalho("O que o vendedor abre na segunda de manhã",
                 "A resposta para o diretor")
cartao(s, MARGEM, y, UTIL, 4.05, tom=CARD, borda=ACENTO, largura_borda=1.5)
px = MARGEM + 0.5
texto(s, "Vendedor:", px, y + 0.3, 1.6, 0.3, 13, MARCA, True)
texto(s, "Quem eu ligo essa semana?", px + 1.7, y + 0.3, 8.0, 0.3, 14.5, BRANCO)
texto(s, "Agente:", px, y + 0.78, 1.6, 0.3, 13, ACENTO, True)
texto(s, "Seus 6 desta semana, em ordem:", px + 1.7, y + 0.78, 8.0, 0.3, 14.5, BRANCO)
itens_chat = [
    ("1. Perfumaria Aurora — score 0,91",
     "Comprava a cada 8 dias, está há 26 sem pedido. Risco de perder.\n"
     "Ofereça: Layali Oud EDP — em estoque."),
    ("2. Drogaria Bella Vita — score 0,87",
     "Comprou o lançamento no mês passado. Alta chance de repetir."),
    ("3. Boutique Essenza — score 0,84",
     "R$ 62 mil no ano. Cliente grande, manter próximo."),
]
cy = y + 1.28
for titulo_, corpo in itens_chat:
    n = corpo.count("\n") + 1
    texto(s, titulo_, px + 0.2, cy, UTIL - 1.4, 0.3, 15, ACENTO_CLR, True)
    texto(s, corpo, px + 0.2, cy + 0.3, UTIL - 1.4, 0.6, 13, CINZA, entrelinha=1.3)
    cy += 0.42 + n * 0.26 + 0.12
texto(s, "Os 200 de maior score da base inteira — e a fatia de cada vendedor "
         "varia de 2 a 10.", MARGEM, 6.42, UTIL, 0.4, 16, ACENTO_CLR, True)
nota(s, "Se perguntarem por que não é a mesma cota para todo mundo: porque a "
        "carteira de um vendedor é mais quente que a do outro. A fila é global, "
        "a capacidade é que é por pessoa. Dar 5 fixos para cada obrigaria a "
        "ligar para cliente frio enquanto sobra cliente quente na mesa do vizinho.")

# ══════════════════════════════════════════════════════════════════════
#  FECHO                                                 slides 38–40
# ══════════════════════════════════════════════════════════════════════

# ── NOVO · os seis motivos da lista ───────────────────────────────────
s, y = cabecalho("A lista não repete a mesma desculpa",
                 "Seis motivos, 200 contatos")
tabela(s, y,
       ["MOTIVO", "CONTATOS", "O QUE O VENDEDOR FAZ"],
       [["Comprou lançamento", "67", "Oferece o resto da linha nova"],
        ["Ciclo curto", "60", "Mantém a cadência, não deixa esfriar"],
        ["Visita converte", "38", "Marca visita em vez de ligar"],
        ["Cliente grande", "19", "Relacionamento, não empurra produto"],
        ["Manutenção", "12", "Contato rápido de rotina"],
        ["Atrasado", "4", "Liga hoje — risco de perder"]],
       larguras=[3.6, 1.7, 6.53], tam=12.5, alt=0.5, destaque=5, cor_destaque=ALERTA)
texto(s, "Se os 200 saíssem com o mesmo motivo, a lista viraria enfeite. "
         "A ordem das regras é que garante isso.",
      MARGEM, 6.5, UTIL, 0.4, 14, ACENTO_CLR, True)

# ── 24 · o antes e o depois ───────────────────────────────────────────
s, y = cabecalho("A resposta completa", "Quais 200?")
texto(s, "ANTES", MARGEM + 4.44, y, 3.4, 0.3, 11, ACENTO, True, spc=140)
texto(s, "DEPOIS", MARGEM + 8.14, y, 3.4, 0.3, 11, ACENTO, True, spc=140)
yy = y + 0.4
for rot, antes, depois, marca in [
        ("Critério",  "Intuição do vendedor", "Score de propensão", None),
        ("Cobertura", "Os 20 que ele lembra", "Os 3.000 avaliados", None),
        ("Ordem",     "Aleatória", "Por probabilidade", None),
        ("Contexto",  "“Faz tempo que não falo\ncom ele”",
                      "Atraso relativo, histórico,\nsugestão", None),
        ("Resultado medido", "20 vendas", "86 vendas", "destaque")]:
    alt = 0.84 if "\n" in antes else 0.6
    cor_b = ACENTO if marca else BORDA
    cartao(s, MARGEM, yy, UTIL, alt, tom=CARD_2 if marca else CARD, borda=cor_b,
           largura_borda=1.5 if marca else 1.25)
    texto(s, rot, MARGEM + 0.34, yy + 0.16, 4.0, alt, 13.5,
          ACENTO if marca else BRANCO, True)
    texto(s, antes, MARGEM + 4.44, yy + 0.16, 3.5, alt, 13.5, CINZA, entrelinha=1.25)
    texto(s, depois, MARGEM + 8.14, yy + 0.16, 3.5, alt, 13.5,
          BRANCO if marca else CINZA, marca is not None, entrelinha=1.25)
    yy += alt + 0.12
bloco(s, 6.28, "Mesmo time. Mesmas 200 ligações. Mais de quatro vezes o resultado.",
      h=0.68, tam=16)

# ── 25 · o arco de três noites ────────────────────────────────────────
s, y = cabecalho("De onde a gente veio", "O arco de três noites")
mono(s, y + 0.05,
     "CSV sujo numa pasta                    segunda\n"
     "         ↓\n"
     "Bronze · Silver · Gold                   terça\n"
     "         ↓\n"
     "Features · Modelo · MLflow              quarta\n"
     "         ↓\n"
     "\"Ligue para esses 200\"                  quarta", tam=15, h=3.55)
texto(s, "E vocês construíram junto. Não assistiram.",
      MARGEM, 6.32, UTIL, 0.45, 18, ACENTO_CLR, True, spc=-15)

# ── 26 · a frase da noite ─────────────────────────────────────────────
impacto("Dashboard descreve o passado.\n"
        "Modelo prevê o futuro.\n"
        "Agente diz o que fazer na segunda de manhã.")

# ══════════════════════════════════════════════════════════════════════
import os
saida = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                     "aula-03-ciencia-de-dados.pptx")
prs.save(saida)
print(f"{len(prs.slides._sldIdLst)} slides → {saida}")
