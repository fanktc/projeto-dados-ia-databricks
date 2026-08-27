#!/usr/bin/env python3
"""Slides da Noite 4 — "E quem não escreve SQL?" · apps e agentes.

Slide também é código. A noite inteira responde UMA pergunta:

    "O pipeline roda, o modelo escolhe os 200. Como o diretor vê isso —
     e como o vendedor devolve o que aconteceu?"

Três prompts ao vivo: o Genie da direção → o app → o retorno da ligação. Os
divisores e os slides de conceito são o que se fala ENQUANTO o Claude Code
trabalha — e nesta noite há uma janela de quase quatro minutos (o primeiro
deploy do app), que os slides 29 a 33 existem para preencher.

Todo número aqui foi medido no workspace em 27/08, com o dataset de seed 42.

Formato 16:9 (13,333 x 7,5 pol). Mesma paleta e mesmas primitivas da noite 3.

Uso:  python gerar_slides.py        (precisa de python-pptx)
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── paleta ────────────────────────────────────────────────────────────
FUNDO      = RGBColor(0x07, 0x0C, 0x16)
CARD       = RGBColor(0x0D, 0x15, 0x22)
CARD_2     = RGBColor(0x14, 0x23, 0x37)
BORDA      = RGBColor(0x1F, 0x3A, 0x5A)
ACENTO     = RGBColor(0x3B, 0x9D, 0xF5)
ACENTO_CLR = RGBColor(0x7C, 0xC8, 0xFF)
ALERTA     = RGBColor(0xFF, 0x6B, 0x5A)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA      = RGBColor(0x9B, 0xAA, 0xC0)
MARCA      = RGBColor(0x56, 0x6A, 0x85)

SERIF = "Arial"
MONO  = "Consolas"

LARGURA, ALTURA = Inches(13.333), Inches(7.5)
MARGEM = 0.75
UTIL   = 13.333 - 2 * MARGEM      # 11,833

prs = Presentation()
prs.slide_width, prs.slide_height = LARGURA, ALTURA
BRANCA = prs.slide_layouts[6]


# ── primitivas ────────────────────────────────────────────────────────
def nova(cor=FUNDO):
    s = prs.slides.add_slide(BRANCA)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = cor
    return s


def texto(s, txt, x, y, w, h, tam, cor, negrito=False, fonte=SERIF,
          alinha=PP_ALIGN.LEFT, entrelinha=1.0, spc=None):
    cx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, linha in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = entrelinha
        r = p.add_run()
        r.text = linha
        r.font.name = fonte
        r.font.size = Pt(tam)
        r.font.bold = negrito
        r.font.color.rgb = cor
        if spc is not None:                     # letter-spacing, em centipontos
            r.font._rPr.set("spc", str(int(spc)))
    return cx


def cartao(s, x, y, w, h, tom=CARD, borda=BORDA, largura_borda=1.25):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = tom
    sh.line.color.rgb = borda
    sh.line.width = Pt(largura_borda)
    sh.shadow.inherit = False
    return sh


def barra(s, cor=ACENTO):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            Inches(0.22), ALTURA)
    sh.fill.solid()
    sh.fill.fore_color.rgb = cor
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def assinatura(s):
    texto(s, "JORNADA DE DADOS", MARGEM, 6.98, 6.0, 0.3, 9, MARCA, True, spc=120)


def nota(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def grade(y, h, n, folga=0.26):
    w = (UTIL - folga * (n - 1)) / n
    return [(MARGEM + i * (w + folga), y, w, h) for i in range(n)]


# ── arquétipos ────────────────────────────────────────────────────────
def cabecalho(kicker, titulo, sub=None, tam=36):
    """Kicker + título + subtítulo. Devolve (slide, y livre)."""
    s = nova()
    texto(s, kicker.upper(), MARGEM, 0.5, UTIL, 0.28, 11.5, ACENTO, True, spc=180)
    n = titulo.count("\n") + 1
    texto(s, titulo, MARGEM, 0.88, UTIL, 0.62 * n, tam, BRANCO, True,
          entrelinha=0.95, spc=-30)
    y = 0.88 + n * (tam / 72 * 1.06) + 0.24
    if sub:
        texto(s, sub, MARGEM, y, UTIL, 0.4, 14, CINZA, entrelinha=1.15)
        y += 0.34 * (sub.count("\n") + 1) + 0.22
    assinatura(s)
    return s, y


def divisor(kicker, titulo, sub):
    s = nova()
    barra(s)
    texto(s, kicker.upper(), 1.15, 2.55, 10.0, 0.32, 13, ACENTO, True, spc=200)
    n = titulo.count("\n") + 1
    texto(s, titulo, 1.15, 3.05, 11.3, 1.1 * n, 60, BRANCO, True,
          entrelinha=0.95, spc=-45)
    texto(s, sub, 1.15, 3.05 + n * 0.92 + 0.42, 10.6, 0.5, 18, CINZA)
    return s


def impacto(linhas):
    s = nova(ACENTO)
    texto(s, linhas, 1.3, 2.55, 10.73, 2.6, 34, BRANCO, True,
          alinha=PP_ALIGN.CENTER, entrelinha=1.35, spc=-20)
    return s


def bloco(s, y, txt, h=0.78, tom=CARD_2, borda=ACENTO, cor=BRANCO, tam=15,
          negrito=True, alinha=PP_ALIGN.LEFT):
    """Faixa de destaque, largura total."""
    cartao(s, MARGEM, y, UTIL, h, tom=tom, borda=borda, largura_borda=1.5)
    n = txt.count("\n") + 1
    texto(s, txt, MARGEM + 0.42, y + (h - n * tam / 72 * 1.28) / 2 + 0.02,
          UTIL - 0.84, h, tam, cor, negrito, alinha=alinha, entrelinha=1.28)


def codigo(s, y, txt, h=None, tam=12.5):
    n = txt.count("\n") + 1
    h = h or n * tam / 72 * 1.42 + 0.5
    cartao(s, MARGEM, y, UTIL, h, tom=CARD, borda=ACENTO, largura_borda=1.25)
    texto(s, txt, MARGEM + 0.4, y + 0.25, UTIL - 0.8, h, tam, ACENTO_CLR,
          fonte=MONO, entrelinha=1.42)
    return h


def mono(s, y, txt, h=None, tam=14, cor=ACENTO_CLR, recuo=None):
    """Diagrama em monoespaçada. recuo=x centraliza o BLOCO (não cada linha),
    que é o que mantém setas e conectores no lugar."""
    n = txt.count("\n") + 1
    h = h or n * tam / 72 * 1.5 + 0.55
    cartao(s, MARGEM, y, UTIL, h)
    if recuo is None:
        largura = max(len(l) for l in txt.split("\n")) * tam / 72 * 0.6
        recuo = max((UTIL - largura) / 2, 0.5)
    texto(s, txt, MARGEM + recuo, y + 0.28, UTIL - recuo, h, tam, cor, True,
          fonte=MONO, entrelinha=1.5)
    return h


def cartoes(s, y, h, itens, tam_rot=17, tam_corpo=12.5, mono_rot=False):
    """itens: (rotulo, corpo) ou (rotulo, corpo, 'alerta'|'destaque')."""
    for (x, yy, w, hh), item in zip(grade(y, h, len(itens)), itens):
        marca = item[2] if len(item) > 2 else None
        cor_b = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BORDA)
        tom   = CARD_2 if marca else CARD
        cartao(s, x, yy, w, hh, tom=tom, borda=cor_b,
               largura_borda=1.5 if marca else 1.25)
        cor_r = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BRANCO)
        nr = item[0].count("\n") + 1
        texto(s, item[0], x + 0.28, yy + 0.26, w - 0.56, 0.9, tam_rot, cor_r,
              True, fonte=MONO if mono_rot else SERIF, entrelinha=1.05)
        texto(s, item[1], x + 0.28, yy + 0.26 + nr * tam_rot / 72 * 1.22 + 0.16,
              w - 0.56, hh, tam_corpo, CINZA, entrelinha=1.28)


def faixas(s, y, itens, alt=0.62, folga=0.14, larg_esq=4.6, tam=13.5,
           cabecalho_=None):
    """Tabela de duas colunas em faixas. itens: (esq, dir) ou (esq, dir, marca)."""
    if cabecalho_:
        texto(s, cabecalho_[0], MARGEM + 0.34, y, larg_esq, 0.3, 11, ACENTO, True, spc=140)
        texto(s, cabecalho_[1], MARGEM + 0.34 + larg_esq, y, UTIL - larg_esq - 0.7,
              0.3, 11, ACENTO, True, spc=140)
        y += 0.42
    for item in itens:
        marca = item[2] if len(item) > 2 else None
        cor_b = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BORDA)
        tom   = CARD_2 if marca else CARD
        cartao(s, MARGEM, y, UTIL, alt, tom=tom, borda=cor_b,
               largura_borda=1.5 if marca else 1.25)
        cor_e = ALERTA if marca == "alerta" else (ACENTO if marca == "destaque" else BRANCO)
        texto(s, item[0], MARGEM + 0.34, y + (alt - tam / 72 * 1.3) / 2,
              larg_esq - 0.2, alt, tam, cor_e, True)
        texto(s, item[1], MARGEM + 0.34 + larg_esq, y + (alt - tam / 72 * 1.3) / 2,
              UTIL - larg_esq - 0.7, alt, tam, CINZA)
        y += alt + folga
    return y


def fluxo(s, y, etapas, h=0.95, quebrado=False, tam=13):
    """Caixas ligadas por seta. quebrado: a última vem em alerta, com ✗."""
    n = len(etapas)
    seta = 0.42
    w = (UTIL - seta * (n - 1)) / n
    for i, e in enumerate(etapas):
        x = MARGEM + i * (w + seta)
        ultimo = i == n - 1
        alvo = ultimo and quebrado
        destaque = ultimo and not quebrado
        cartao(s, x, y, w, h,
               tom=CARD_2 if (alvo or destaque) else CARD,
               borda=ALERTA if alvo else (ACENTO if destaque else BORDA),
               largura_borda=1.5 if (alvo or destaque) else 1.25)
        texto(s, e, x + 0.16, y + (h - tam / 72 * 1.3 * (e.count("\n") + 1)) / 2,
              w - 0.32, h, tam, ALERTA if alvo else (ACENTO_CLR if destaque else BRANCO),
              True, fonte=MONO, alinha=PP_ALIGN.CENTER, entrelinha=1.3)
        if not ultimo:
            texto(s, "✗" if (quebrado and i == n - 2) else "→",
                  x + w, y + h / 2 - 0.19, seta, 0.4, 17,
                  ALERTA if (quebrado and i == n - 2) else ACENTO, True,
                  alinha=PP_ALIGN.CENTER)



# ── arquétipo 6: tabela de dados ──────────────────────────────────────
def tabela(s, y, colunas, linhas, larguras=None, tam=13, alt=0.58,
           destaque=None, cor_destaque=None):
    """Uma tabela com dados reais. colunas: lista de rótulos. linhas: lista de
    listas. destaque: índice da linha a realçar."""
    n = len(colunas)
    larguras = larguras or [UTIL / n] * n
    xs, acum = [], MARGEM
    for w in larguras:
        xs.append(acum); acum += w

    for x, w, rot in zip(xs, larguras, colunas):
        texto(s, rot, x + 0.22, y, w - 0.3, 0.3, 11, ACENTO, True, spc=120)
    yy = y + 0.4
    for i, linha in enumerate(linhas):
        marcada = destaque is not None and i == destaque
        cor = cor_destaque or ACENTO
        cartao(s, MARGEM, yy, UTIL, alt,
               tom=CARD_2 if marcada else CARD,
               borda=cor if marcada else BORDA,
               largura_borda=1.5 if marcada else 1.25)
        for x, w, val in zip(xs, larguras, linha):
            negrito = marcada or x == xs[0]
            cor_txt = (cor if marcada and x != xs[0] else
                       (BRANCO if x == xs[0] else CINZA))
            texto(s, str(val), x + 0.22, yy + (alt - tam / 72 * 1.3) / 2,
                  w - 0.3, alt, tam, cor_txt, negrito,
                  fonte=MONO if x != xs[0] else SERIF)
        yy += alt + 0.1
    return yy


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 1 · A PERGUNTA                                    slides 1–6
# ══════════════════════════════════════════════════════════════════════

# ── 1 · capa ──────────────────────────────────────────────────────────
s = nova()
barra(s)
texto(s, "IMERSÃO AO VIVO · NOITE 4 DE 4 · 27/08 · 19H30",
      1.15, 2.25, 11.0, 0.32, 13, ACENTO, True, spc=200)
texto(s, "E quem não\nescreve SQL?", 1.1, 2.62, 11.3, 2.4, 78, BRANCO, True,
      entrelinha=0.92, spc=-60)
texto(s, "Apps, agentes e o caminho de volta",
      1.15, 5.15, 11.0, 0.5, 25, ACENTO_CLR, True, spc=-15)
texto(s, "Rota do Perfume · Databricks Apps + Genie + Claude Code",
      1.15, 6.55, 11.0, 0.35, 13, MARCA)
nota(s, "Abra com a URL do app FECHADA. A sala precisa sentir a falta antes "
        "de ver a solução.")

# ── 2 · a ligação de novo ─────────────────────────────────────────────
s, y = cabecalho("Ontem à noite, depois da aula", "O diretor ligou de novo")
cartao(s, MARGEM, y, UTIL, 3.5, tom=CARD, borda=ACENTO, largura_borda=1.5)
texto(s, "“A lista ficou ótima. Os 200 nomes, com o motivo escrito. Perfeito.",
      MARGEM + 0.55, y + 0.4, UTIL - 1.1, 0.5, 16, BRANCO, entrelinha=1.3)
texto(s, "Só que ela está numa tela que só você abre. Eu pedi para o meu\n"
         "gerente ver e ele me mandou um print. Print, Luciano.",
      MARGEM + 0.55, y + 1.05, UTIL - 1.1, 1.0, 16, CINZA, entrelinha=1.32)
texto(s, "Me manda um link. E eu quero marcar quem já foi contatado.”",
      MARGEM + 0.55, y + 2.5, UTIL - 1.1, 0.6, 22, ACENTO_CLR, True, spc=-20)
bloco(s, 6.1, "Essa é a aula de hoje. Duas frases, dois prompts.", h=0.62, tam=16)

# ── 3 · quem consegue usar o que construímos ──────────────────────────
s, y = cabecalho("Três noites de trabalho",
                 "O que era preciso saber para usar")
faixas(s, y, [
    ("Noite 1 · o dado no catálogo", "SQL"),
    ("Noite 2 · o pipeline, o dashboard e o Genie", "Python e PySpark"),
    ("Noite 3 · o modelo e a fila dos 200", "ML e SQL", "alerta"),
    ("Noite 4 · o app e o Genie da direção", "nada — e é a aula de hoje", "destaque"),
], alt=0.72, larg_esq=5.6, cabecalho_=("O QUE FICOU DE PÉ", "O QUE ERA PRECISO SABER"))
bloco(s, 5.95, "A empresa tem 3.000 clientes, 42 vendedores e um diretor.\n"
                "Nenhum deles escreve SQL.", h=1.0, tam=16, borda=ALERTA, cor=BRANCO)

# ── 4 · a frase da noite ──────────────────────────────────────────────
impacto("Um dado que só o time de dados\n"
        "consegue abrir é um dado que\n"
        "não existe para a empresa.")

# ── 5 · o que a noite entrega ─────────────────────────────────────────
s, y = cabecalho("A promessa de hoje", "O projeto ganha uma URL",
                 "E, no último prompt, ganha um caminho de volta — que é a parte "
                 "que quase todo projeto de dados esquece.")
cartoes(s, y + 0.1, 2.5, [
    ("1 · Perguntar",
     "Um Genie feito para a direção: a fila, o que ela vale, o que aconteceu "
     "depois. Fala em ligações, nunca em AUC."),
    ("2 · Ver",
     "Um Databricks App com os 200 na tela, filtráveis por vendedor, e o Genie "
     "embutido numa aba."),
    ("3 · Devolver",
     "Quatro botões. O vendedor marca como foi a ligação, e isso vira uma linha "
     "na gold.", "destaque"),
])
bloco(s, 5.9, "Nenhuma tabela de análise nova. Tudo que a noite mostra já "
               "estava pronto ontem.", h=0.72, tam=15.5)

# ── 6 · o plano ───────────────────────────────────────────────────────
s, y = cabecalho("O formato de sempre", "Três prompts, três deploys")
h = mono(s, y + 0.05,
     "prompt 1   + genie_direcao      o Genie da direção · gold.retorno_ligacao\n"
     "prompt 2   + app                a fila dos 200 na tela, com o Genie dentro\n"
     "prompt 3   + POST /api/retorno  o resultado da ligação volta para a gold")
y2 = y + h + 0.3
faixas(s, y2, [
    ("bundle deploy", "~20s · o Genie e a tabela entram no bundle da terça"),
    ("apps deploy (o primeiro)", "3m44s · o app tem compute próprio, e ele nasce agora", "destaque"),
    ("apps deploy (o segundo)", "1m04s · daí em diante é rápido"),
], alt=0.66, larg_esq=5.0)
bloco(s, 6.25, "O job vai de 15 para 16 tarefas. O app é o único artefato da "
               "imersão com ciclo de deploy próprio.", h=0.7, tam=15)


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 2 · A RECAPITULAÇÃO EM VISÃO DE NEGÓCIO           slides 7–14
# ══════════════════════════════════════════════════════════════════════

# ── 7 · divisor ───────────────────────────────────────────────────────
divisor("Antes de abrir qualquer código",
        "O que a gente\nconstruiu",
        "Três noites contadas na língua de quem paga a conta — sem uma palavra "
        "de arquitetura.")

# ── 8 · noite 1 ───────────────────────────────────────────────────────
s, y = cabecalho("Noite 1 · segunda", "“Qual foi a receita?”",
                 "A pergunta mais simples que existe. E ela quebrou.")
cartoes(s, y + 0.1, 2.25, [
    ("O que existia",
     "Dez arquivos exportados do ERP e do CRM. 28.729 pedidos, 3.040 clientes, "
     "197.724 itens."),
    ("O que aconteceu",
     "A soma da receita deu errado. Data em dois formatos, cliente duplicado, "
     "devolução somando como venda."),
    ("O que o negócio aprendeu",
     "O número que a diretoria vê depende de decisões que ninguém documentou.",
     "destaque"),
])
bloco(s, 5.65, "Fim da noite 1: o dado estava no lugar. Confiável, ainda não.",
      h=0.7, tam=16)

# ── 9 · noite 2 · a limpeza ───────────────────────────────────────────
s, y = cabecalho("Noite 2 · terça", "“Como eu não repito isso todo mês?”",
                 "Seis prompts, seis deploys. O trabalho manual virou pipeline.")
tabela(s, y + 0.05,
       ["O QUE ESTAVA ERRADO", "QUANTO", "O QUE ISSO CUSTAVA"],
       [["Datas em dd/mm/aaaa", "3.443", "relatório de mês errado"],
        ["Clientes duplicados", "40", "carteira contada duas vezes"],
        ["Cancelados com valor zero", "957", "pedido que nunca existiu"],
        ["Devoluções negativas", "2.327", "receita inflada"]],
       larguras=[5.2, 1.9, 4.733], alt=0.56)
bloco(s, 5.7, "Receita que aparecia e não existia:  R$ 3.586.620,37",
      h=0.78, tam=19, borda=ALERTA, cor=BRANCO, alinha=PP_ALIGN.CENTER)
nota(s, "R$ 105.890.448,42 no item cru contra R$ 102.303.828,05 na gold.")

# ── 10 · noite 2 · as portas ──────────────────────────────────────────
s, y = cabecalho("Noite 2 · terça", "E a diretoria ganhou duas portas")
cartoes(s, y + 0.15, 2.45, [
    ("O dashboard",
     "Receita por mês, margem por categoria, ranking de marcas. Atualiza "
     "sozinho todo dia, porque o pipeline roda sozinho."),
    ("O Genie comercial",
     "Pergunta em português sobre a gold. Sabe que o pico do setor é o mês "
     "ANTERIOR à data comemorativa — porque alguém escreveu isso."),
    ("Os testes que quebram",
     "Nove conferências. Se o número não fechar, o job PARA. Melhor o "
     "dashboard com o dado de ontem que com o dado errado de hoje.", "destaque"),
])
bloco(s, 5.95, "Fim da noite 2: o número certo, todo dia, sem ninguém tocar.",
      h=0.7, tam=16)

# ── 11 · noite 3 · o modelo ───────────────────────────────────────────
s, y = cabecalho("Noite 3 · quarta", "“Quais 200 eu ligo?”",
                 "3.000 clientes na base, 200 ligações boas por semana.")
tabela(s, y + 0.05,
       ["A ESTRATÉGIA", "DOS 200, QUANTOS COMPRAM"],
       [["Ligar aleatório", "20"],
        ["Ligar para quem sumiu há mais tempo", "0"],
        ["Ligar para os maiores clientes", "44"],
        ["Ligar para os 200 de maior score", "86"]],
       larguras=[7.4, 4.433], alt=0.6, destaque=3)
bloco(s, 5.85, "4,25× mais pedidos com a mesma equipe, no mesmo horário, "
               "com o mesmo telefone.", h=0.78, tam=17, alinha=PP_ALIGN.CENTER)
nota(s, "A segunda linha é o momento da noite 3: a intuição comercial não "
        "estava imprecisa, estava invertida.")

# ── 12 · noite 3 · o que a fila vale ──────────────────────────────────
s, y = cabecalho("Noite 3 · quarta", "E a lista tem preço")
cartoes(s, y + 0.15, 2.35, [
    ("200 contatos",
     "Os 200 maiores scores da base inteira, divididos pela carteira de 35 "
     "vendedores. Não é cota igual: é capacidade por pessoa."),
    ("R$ 582.799,50",
     "A soma de score × ticket médio de cada cliente da fila. É estimativa, "
     "e é o número que vai para a reunião de segunda.", "destaque"),
    ("Um motivo por linha",
     "“Compra a cada 8 dias e está há 26 sem pedido.” Score não é decisão — "
     "0,84 não é uma ação."),
])
bloco(s, 5.85, "Fim da noite 3: o pipeline parou de contar o passado e passou "
               "a dizer o que fazer.", h=0.7, tam=16)

# ── 13 · o arco numa query ────────────────────────────────────────────
s, y = cabecalho("As três noites", "Em quatro números")
tabela(s, y + 0.1,
       ["ETAPA", "NÚMERO", "O QUE É"],
       [["noite 1 · o dado", "28.729", "pedidos que chegaram do ERP"],
        ["noite 2 · o pipeline", "191.080", "linhas em gold.fato_vendas"],
        ["noite 3 · a decisão", "200", "contatos na fila desta semana"],
        ["noite 4 · o retorno", "?", "ligações que o time registrou"]],
       larguras=[4.2, 2.4, 5.233], alt=0.62, destaque=3, cor_destaque=ALERTA)
bloco(s, 6.0, "A última linha é a de hoje. Ela vale zero até alguém clicar "
               "num botão.", h=0.7, tam=16, borda=ALERTA)
nota(s, "Rode esta query no fecho da noite — está no QUERIES.md.")

# ── 14 · o que falta ──────────────────────────────────────────────────
s, y = cabecalho("O que as três noites não resolveram", "O caminho de volta")
h = mono(s, y + 0.25,
     "   pipeline  →  score  →  fila  →  ligação  →  ???\n"
     "                  ↑                              │\n"
     "                  └──────────────────────────────┘\n"
     "                     é este pedaço que falta", tam=15)
bloco(s, y + h + 0.55,
      "O modelo diz que 86 dos 200 vão comprar.\nNinguém aqui sabe se ele acertou.",
      h=1.05, tam=18, borda=ALERTA, alinha=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 3 · AS TRÊS PORTAS                                slides 15–18
# ══════════════════════════════════════════════════════════════════════

# ── 15 · divisor ──────────────────────────────────────────────────────
divisor("A decisão de desenho da noite",
        "Três portas,\numa tabela",
        "Dashboard, Genie e app leem exatamente o mesmo gold.fila_semanal. "
        "Não é redundância — é audiência.")

# ── 16 · quando usar cada uma ─────────────────────────────────────────
s, y = cabecalho("A pergunta que a sala faz", "Por que três, se o dado é o mesmo?")
tabela(s, y + 0.05,
       ["PORTA", "PARA QUEM", "O QUE SÓ ELA FAZ"],
       [["Dashboard", "quem acompanha número recorrente", "agenda e alerta, sem código"],
        ["Genie", "quem tem pergunta não prevista", "responde o que não estava na tela"],
        ["App", "quem trabalha a lista todo dia", "interação e escrita de volta"]],
       larguras=[2.4, 4.6, 4.833], alt=0.66, destaque=2)
bloco(s, 5.55, "Genie responde. App registra.\n"
               "A diferença não é a tecnologia — é a direção do dado.",
      h=1.05, tam=18, alinha=PP_ALIGN.CENTER)

# ── 17 · o que é databricks apps ──────────────────────────────────────
s, y = cabecalho("A peça nova de hoje", "Databricks Apps",
                 "Uma aplicação web que roda DENTRO do workspace, com identidade "
                 "própria e acesso governado pelo Unity Catalog.")
cartoes(s, y + 0.1, 2.5, [
    ("Não é servidor seu",
     "Sem VM, sem Docker, sem domínio para comprar. O deploy é um comando e "
     "sai uma URL do workspace."),
    ("Login é o do Databricks",
     "Quem não tem acesso ao workspace não abre. O app sabe QUEM está olhando "
     "— e é assim que o retorno grava quem clicou."),
    ("Tem um usuário próprio",
     "Um service principal criado junto com o app. E ele nasce sem permissão "
     "nenhuma no catálogo.", "alerta"),
])
bloco(s, 5.9, "O último cartão é o erro nº 1 de quem sobe o primeiro app. "
               "Voltamos nele no slide 29.", h=0.72, tam=15)

# ── 18 · por que dois genies ──────────────────────────────────────────
s, y = cabecalho("A outra peça, que já é velha conhecida",
                 "Por que um SEGUNDO Genie?")
y = faixas(s, y + 0.05, [
    ("Rota do Perfume · Comercial", "12 fontes · da noite 2 · para o vendedor perguntar qualquer coisa da operação"),
    ("Rota do Perfume · Direção", "7 fontes · nasce hoje · para responder UMA decisão: quem ligar e quanto vale", "destaque"),
], alt=0.9, larg_esq=5.2, cabecalho_=("O SPACE", "O RECORTE"))
bloco(s, y + 0.28,
      "Genie não é um por empresa. É um por audiência.\n"
      "O que muda entre eles não é o dado — é o recorte e a instrução.",
      h=1.1, tam=17, alinha=PP_ALIGN.CENTER)
nota(s, "Se as duas audiências moram no mesmo space, as instruções brigam: o "
        "que serve para o vendedor vira ruído para o diretor.")


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 4 · PROMPT 1 · O GENIE DA DIREÇÃO                 slides 19–24
# ══════════════════════════════════════════════════════════════════════

# ── 19 · divisor ──────────────────────────────────────────────────────
divisor("Prompt 1 de 3 · deploy nº 1",
        "O Genie\nda direção",
        "O space curado como código, mais a tabela que vai receber a resposta "
        "do time — e que nasce vazia.")

# ── 20 · a instrução é o produto ──────────────────────────────────────
s, y = cabecalho("O que faz um Genie responder melhor que o outro",
                 "A instrução é o produto",
                 "O modelo é o mesmo. O dado é o mesmo. O que muda são vinte "
                 "linhas escritas por alguém que entende o negócio.")
faixas(s, y + 0.05, [
    ("“A fila é global, não é cota”", "quem tem carteira quente recebe mais contatos, e isso está certo"),
    ("“Receita esperada é ESTIMATIVA”", "SUM(score × ticket_medio) nunca é receita realizada"),
    ("“NUNCA cite AUC”", "AUC é métrica de quem treina. A métrica da direção é lift_top200", "destaque"),
    ("“Se for zero, diga que é zero”", "retorno_ligacao começa vazia — não invente, e não use a fila no lugar"),
], alt=0.66, larg_esq=4.8)
bloco(s, 6.2, "E tudo isso mora no Git: com histórico, revisão e rollback.",
      h=0.66, tam=15.5)

# ── 21 · a tabela que nasce vazia ─────────────────────────────────────
s, y = cabecalho("A outra entrega do prompt 1", "A tabela que nasce vazia")
h = codigo(s, y + 0.05,
     "CREATE TABLE IF NOT EXISTS gold.retorno_ligacao (\n"
     "  cliente_id      INT,      vendedor    STRING,\n"
     "  status          STRING,   comentario  STRING,\n"
     "  registrado_em   TIMESTAMP,\n"
     "  registrado_por  STRING,   _referencia DATE\n"
     ")", tam=13)
y2 = y + h + 0.3
cartoes(s, y2, 1.85, [
    ("IF NOT EXISTS",
     "Toda tabela da gold é recriada a cada execução. Esta não pode ser: o dado "
     "dela não vem do pipeline, vem do time.", "destaque"),
    ("status é um conjunto fechado",
     "vendeu · vai_pensar · sem_interesse · nao_atendeu. Sem isso, em três "
     "semanas a coluna tem “Vendeu”, “vendido” e “VENDEU”."),
])
bloco(s, 6.3, "É a única tabela do projeto que o pipeline cria e nunca recria.",
      h=0.66, tam=15.5)

# ── 22 · prompt 1 rodando ─────────────────────────────────────────────
s, y = cabecalho("Agora sim", "Prompt 1 · rodando",
                 "Enquanto ele trabalha, você fala — e isto leva 20 segundos, "
                 "então fale rápido.")
h = mono(s, y + 0.1,
     "resources/direcao.geniespace.json      as 7 fontes e as instruções\n"
     "resources/genie-direcao.genie_space.yml  o recurso no bundle\n"
     "src/gold/11-retorno-ligacao.sql        a tabela do caminho de volta\n"
     "resources/pipeline.job.yml             +1 tarefa  →  16", tam=13)
bloco(s, y + h + 0.45,
      "NÃO use --auto-approve.\n"
      "Se o deploy pedir para apagar o dashboard ou o Genie comercial, pare.",
      h=1.05, tam=16, borda=ALERTA, alinha=PP_ALIGN.CENTER)

# ── 23 · as três perguntas ────────────────────────────────────────────
s, y = cabecalho("A conferência", "Três perguntas, na frente da sala")
tabela(s, y + 0.05,
       ["A PERGUNTA", "O QUE TEM QUE APARECER"],
       [["Quanto vale a fila desta semana?", "R$ 582.799,50 — e a palavra estimativa"],
        ["Quantas ligações já foram registradas?", "zero, e a frase de que ninguém registrou"],
        ["O modelo é bom?", "4,25× ou 86 de 200. Nunca AUC"]],
       larguras=[5.8, 6.033], alt=0.66, destaque=2)
bloco(s, 5.65, "Use “Show generated code” em toda resposta.\n"
               "É o hábito que separa quem USA Genie de quem CONFIA em Genie.",
      h=1.05, tam=17, alinha=PP_ALIGN.CENTER)

# ── 24 · o que acabou de acontecer ────────────────────────────────────
s, y = cabecalho("Deploy nº 1 · pronto", "O que existe agora que não existia")
cartoes(s, y + 0.2, 2.4, [
    ("Um segundo Genie",
     "Sete fontes, instrução própria, no Git. O comercial continua de pé, "
     "intacto."),
    ("Uma tabela vazia",
     "gold.retorno_ligacao, com COMMENT em toda coluna — a auditoria da noite 2 "
     "vale para ela também."),
    ("16 tarefas no job",
     "A noite 4 acrescentou uma, e ela não calcula nada: só garante que a "
     "tabela existe.", "destaque"),
])
bloco(s, 5.9, "E o diretor ainda não tem onde clicar. É o próximo prompt.",
      h=0.7, tam=16)


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 5 · PROMPT 2 · O APP                              slides 25–33
# ══════════════════════════════════════════════════════════════════════

# ── 25 · divisor ──────────────────────────────────────────────────────
divisor("Prompt 2 de 3 · deploy nº 2",
        "Uma URL",
        "A fila dos 200 na tela, filtrável por vendedor, com o Genie do prompt "
        "anterior embutido numa aba.")

# ── 26 · a tela que ele tem hoje ──────────────────────────────────────
s, y = cabecalho("Antes de escrever qualquer coisa", "O que o diretor tem hoje")
h = codigo(s, y + 0.05,
     "SELECT vendedor, ordem, razao_social, ROUND(score,2) AS nota,\n"
     "       motivo, sugestao\n"
     "FROM   lakehouse_rotaperfume.gold.fila_semanal\n"
     "ORDER  BY score DESC;", tam=13)
bloco(s, y + h + 0.4,
      "Está certíssimo. Agora imagine mandar isso para o diretor toda segunda "
      "de manhã.\nEle vai pedir para filtrar por vendedor. Depois, para marcar "
      "quem já foi contatado.\nNa terceira semana, ele volta a ligar pela intuição.",
      h=1.5, tam=15.5, negrito=False, cor=CINZA)

# ── 27 · anatomia do app ──────────────────────────────────────────────
s, y = cabecalho("O que o prompt vai montar", "Três telas, quatro queries")
h = mono(s, y + 0.05,
     "config/queries/kpis_semana.sql     os quatro números do topo\n"
     "config/queries/vendedores.sql      o filtro\n"
     "config/queries/fila.sql            os 200, com o retorno de cada um\n"
     "config/queries/acompanhamento.sql  o desfecho por vendedor\n"
     "\n"
     "client/src/pages/  A semana  ·  Acompanhamento  ·  Perguntar", tam=13)
bloco(s, y + h + 0.45,
      "Nenhuma query dentro do React. SQL em arquivo .sql, interface em "
      "arquivo de interface —\ne o nome do arquivo é a chave que a tela pede.",
      h=1.1, tam=15.5)

# ── 28 · os quatro cartões ────────────────────────────────────────────
s, y = cabecalho("A tela principal", "Os quatro números do diretor",
                 "Medidos hoje, no workspace. É o que ele vê ao abrir o link.")
cartoes(s, y + 0.1, 2.4, [
    ("200\ncontatos", "em 35 vendedores"),
    ("R$ 582.799,50", "receita esperada — soma de score × ticket médio"),
    ("43%", "conversão prevista, contra 10,1% ligando às cegas", "destaque"),
    ("0", "ligações registradas. Por enquanto.", "alerta"),
], tam_rot=22, tam_corpo=12)
bloco(s, 5.85, "O terceiro cartão é o projeto inteiro numa linha — e agora "
               "está numa página que o diretor abre sozinho.", h=0.72, tam=15.5)

# ── 29 · o app é um usuário do UC ─────────────────────────────────────
s, y = cabecalho("O erro nº 1 de quem sobe o primeiro app",
                 "O app é um usuário do Unity Catalog")
h = codigo(s, y + 0.05,
     "# no databricks.yml, isto dá acesso ao COMPUTE:\n"
     "  - name: sql-warehouse\n"
     "    sql_warehouse: { id: ..., permission: CAN_USE }\n"
     "\n"
     "# e isto, ao DADO — sem os três, a tela carrega VAZIA:\n"
     "GRANT USE CATALOG ON CATALOG lakehouse_rotaperfume TO `<sp>`;\n"
     "GRANT USE SCHEMA  ON SCHEMA  lakehouse_rotaperfume.gold TO `<sp>`;\n"
     "GRANT SELECT      ON SCHEMA  lakehouse_rotaperfume.gold TO `<sp>`;", tam=12.5)
bloco(s, y + h + 0.35,
      "A tela do erro é a pior possível: ela carrega, não quebra, e mostra "
      "vazio.", h=0.8, tam=16, borda=ALERTA, alinha=PP_ALIGN.CENTER)
nota(s, "O service principal muda a cada app criado. Leia com "
        "`databricks apps get`, nunca copie de outro ambiente.")

# ── 30 · os tipos vêm do catálogo ─────────────────────────────────────
s, y = cabecalho("A recompensa da noite 2, três dias depois",
                 "Os tipos vêm do catálogo")
h = codigo(s, y + 0.05,
     "// gerado por `npm run typegen`, lendo o Unity Catalog:\n"
     "\n"
     "/** Probabilidade de o cliente fazer pedido nos próximos 7 dias. */\n"
     "score: number;\n"
     "/** Por que este cliente está na lista, escrito para o vendedor ler. */\n"
     "motivo: string;", tam=13)
bloco(s, y + h + 0.4,
      "O COMMENT que a auditoria exigiu na terça virou documentação dentro do "
      "editor.\nMetadado não é documentação para humano ler: é interface.",
      h=1.15, tam=16, alinha=PP_ALIGN.CENTER)

# ── 31 · prompt 2 rodando ─────────────────────────────────────────────
s, y = cabecalho("Agora sim", "Prompt 2 · rodando",
                 "Este é o deploy mais longo da imersão inteira. Não é "
                 "travamento — é o compute do app nascendo.")
tabela(s, y + 0.05,
       ["ETAPA", "TEMPO MEDIDO", "O QUE ACONTECE"],
       [["databricks apps init", "~60s", "clona o template e instala dependência"],
        ["typegen + validate", "~6s", "descreve as queries no warehouse"],
        ["apps deploy (1º)", "3m44s", "provisiona compute, instala, faz build"],
        ["apps deploy (2º)", "1m04s", "daí em diante é rápido"]],
       larguras=[4.2, 2.6, 5.033], alt=0.6, destaque=2)
bloco(s, 6.0, "Tenha os slides 29 e 30 prontos para preencher esses "
               "quatro minutos.", h=0.66, tam=15.5)

# ── 32 · a conferência ────────────────────────────────────────────────
s, y = cabecalho("Deploy nº 2 · pronto", "Confira contra o banco, lado a lado")
faixas(s, y + 0.05, [
    ("Contatos · 200 · 35 vendedores", "COUNT(*) e COUNT(DISTINCT vendedor) em fila_semanal"),
    ("Receita esperada · R$ 582.799,50", "SUM(score * ticket_medio)"),
    ("Conversão prevista · 43%", "acertos_top200 / 200, da última versão de modelo_metricas", "destaque"),
    ("Filtro · Débora Souza · 12", "o maior número de contatos da fila"),
], alt=0.66, larg_esq=5.4, cabecalho_=("O QUE A TELA MOSTRA", "A QUERY QUE PROVA"))
bloco(s, 6.2, "Abra o app e o SQL Editor lado a lado. Todo número da tela tem "
               "uma query atrás.", h=0.66, tam=15.5)

# ── 33 · o genie dentro do app ────────────────────────────────────────
s, y = cabecalho("E a terceira aba", "O mesmo Genie, agora dentro do produto")
cartoes(s, y + 0.2, 2.4, [
    ("Uma definição",
     "O space do prompt 1 é o mesmo objeto. O app não tem uma cópia — ele "
     "aponta para o recurso."),
    ("Duas portas",
     "Quem prefere a interface do Databricks abre o space. Quem vive no app "
     "pergunta na aba, sem trocar de tela."),
    ("O SQL sempre à vista",
     "Toda resposta traz a query que a produziu. Resposta de IA sem o caminho "
     "não vai para reunião.", "destaque"),
])
bloco(s, 5.9, "O diretor já tem link e já pergunta. Falta ele conseguir "
               "RESPONDER.", h=0.7, tam=16)


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 6 · PROMPT 3 · O RETORNO                          slides 34–39
# ══════════════════════════════════════════════════════════════════════

# ── 34 · divisor ──────────────────────────────────────────────────────
divisor("Prompt 3 de 3 · deploy nº 3",
        "O caminho\nde volta",
        "Quatro botões. O que o vendedor responde hoje é o rótulo de treino da "
        "semana que vem.")

# ── 35 · o ciclo quebrado ─────────────────────────────────────────────
s, y = cabecalho("O que o projeto ainda não sabe", "O dado sai e não volta")
h = mono(s, y + 0.25,
     "   pipeline  →  score  →  fila  →  ligação  →  ???\n"
     "                  ↑                              │\n"
     "                  └──────────────────────────────┘", tam=15)
bloco(s, y + h + 0.5,
      "É o pedaço mais barato do projeto inteiro — uma tabela e um botão.\n"
      "E é o que separa um relatório de um produto de dados.",
      h=1.2, tam=17, alinha=PP_ALIGN.CENTER)

# ── 36 · o retorno é o rótulo ─────────────────────────────────────────
s, y = cabecalho("Por que isso importa para o modelo",
                 "O retorno de hoje é o rótulo de amanhã")
y = faixas(s, y + 0.05, [
    ("Como a noite 3 treinou", "com comprou_em_7d, calculado do histórico de pedidos"),
    ("O que isso não enxerga", "o cliente que não comprou PORQUE ninguém ligou para ele"),
    ("O que o retorno acrescenta", "ligou e vendeu · ligou e não vendeu · nem ligou — três coisas diferentes", "destaque"),
], alt=0.72, larg_esq=5.0)
bloco(s, y + 0.28,
      "Sem o retorno, o modelo aprende sobre quem comprou.\n"
      "Com ele, aprende sobre quem foi abordado — que é a decisão real.",
      h=1.1, tam=16.5, alinha=PP_ALIGN.CENTER)

# ── 37 · botão é interface, enum é contrato ───────────────────────────
s, y = cabecalho("A decisão de engenharia do prompt 3",
                 "Botão é interface. O enum é o contrato")
h = codigo(s, y + 0.05,
     "const RetornoSchema = z.object({\n"
     "  cliente_id: z.number().int(),\n"
     "  status:     z.enum(['vendeu', 'vai_pensar',\n"
     "                      'sem_interesse', 'nao_atendeu']),\n"
     "  comentario: z.string().max(500).default(''),\n"
     "});   // corpo inválido → 400, sem tocar no warehouse", tam=12.5)
bloco(s, y + h + 0.4,
      "Se o front puder mandar status livre, em três semanas a coluna tem\n"
      "“vendeu”, “Vendeu”, “vendido” e “VENDEU” — e nenhum relatório fecha.",
      h=1.15, tam=16, borda=ALERTA, alinha=PP_ALIGN.CENTER)

# ── 38 · a permissão de escrita ───────────────────────────────────────
s, y = cabecalho("A pergunta para a sala",
                 "Que permissão o app precisa para escrever?")
y = faixas(s, y + 0.15, [
    ("MODIFY no schema gold", "o app pode alterar fato_vendas, os marts, a fila. Tudo", "alerta"),
    ("MODIFY na tabela retorno_ligacao", "o app só pode escrever onde ele é dono do dado", "destaque"),
], alt=0.9, larg_esq=5.6, cabecalho_=("O GRANT", "O QUE ELE PERMITE"))
h = codigo(s, y + 0.18,
     "GRANT MODIFY ON TABLE gold.retorno_ligacao TO `<sp>`;", tam=13.5)
bloco(s, y + 0.18 + h + 0.32,
      "Em TABLE, não em SCHEMA. Permissão é desenho, não burocracia.",
      h=0.72, tam=16, alinha=PP_ALIGN.CENTER)

# ── 39 · o teste do ciclo, passo a passo ──────────────────────────────
s, y = cabecalho("Antes de clicar em nada", "O teste que prova o ciclo inteiro",
                 "Sete passos com o SQL Editor e o app lado a lado. É o que "
                 "transforma “o app grava” em “eu vi gravar”.")
tabela(s, y + 0.05,
       ["PASSO", "ONDE", "O QUE VOCÊ VÊ"],
       [["1 · o estado inicial", "SQL Editor", "SELECT COUNT(*) em retorno_ligacao → 0"],
        ["2 · escolher quem ligar", "SQL Editor", "cliente 2137 · Farmácia Serena · score 0,974"],
        ["3 · o mesmo cliente com LEFT JOIN", "SQL Editor", "retorno e registrado_por vêm NULL"],
        ["4 · o apontamento", "o app", "comentário + Vendeu. O cartão sobe na hora"],
        ["5 · a mesma query de novo", "SQL Editor", "a linha existe, com o SEU e-mail"],
        ["6 · o acompanhamento", "os dois", "o número da tela bate com o do banco"],
        ["7 · a pergunta ao Genie", "Genie", "ele responde o número novo, sem mudar nada"]],
       larguras=[4.0, 2.5, 5.333], alt=0.52, destaque=4)
bloco(s, 6.35, "O passo 5 é o que a sala precisa ver: o NULL de trinta "
               "segundos atrás virou dado.", h=0.66, tam=15.5)
nota(s, "O roteiro completo, com as queries prontas, está em "
        "passo-a-passo/03-retorno.md.")

# ── 40 · o antes e o depois, na mesma query ───────────────────────────
s, y = cabecalho("O passo 3 e o passo 5", "A mesma query, trinta segundos depois")
h = codigo(s, y + 0.05,
     "SELECT f.razao_social, r.status AS retorno, r.registrado_por\n"
     "FROM   gold.fila_semanal f\n"
     "LEFT JOIN gold.retorno_ligacao r ON r.cliente_id = f.cliente_id\n"
     "WHERE  f.cliente_id = 2137;", tam=13)
y2 = y + h + 0.3
cartoes(s, y2, 1.75, [
    ("ANTES do clique",
     "retorno → NULL\nregistrado_por → NULL\n\nO pipeline sabe a quem ligar e "
     "não sabe o que aconteceu."),
    ("DEPOIS do clique",
     "retorno → vendeu\nregistrado_por → o seu e-mail\n\nO caminho de volta "
     "existe, e tem nome e hora.", "destaque"),
])
bloco(s, 6.3, "Não é o app que impressiona. É o NULL que virou dado.",
      h=0.66, tam=16, alinha=PP_ALIGN.CENTER)

# ── 41 · o momento da noite ───────────────────────────────────────────
s, y = cabecalho("Deploy nº 3 · pronto", "Agora clique — e volte para o SQL")
tabela(s, y + 0.05,
       ["A SEQUÊNCIA", "O QUE A SALA VÊ"],
       [["1 · Farmácia Serena, Goiânia", "o primeiro da fila, chance de 97%"],
        ["2 · comentário e “Vendeu”", "o cartão vai de 0 para 1"],
        ["3 · SELECT na retorno_ligacao", "a linha está lá, com o SEU e-mail"],
        ["4 · a mesma pergunta ao Genie", "responde 1 e 1 — e há 20 min respondia zero"]],
       larguras=[5.4, 6.433], alt=0.62, destaque=3)
bloco(s, 5.9, "Nenhuma linha do Genie mudou. Mudou o dado embaixo dele.",
      h=0.72, tam=17, alinha=PP_ALIGN.CENTER)
nota(s, "É o melhor momento da noite. Não corra: faça os quatro passos "
        "devagar, com a sala olhando a tela.")


# ══════════════════════════════════════════════════════════════════════
#  BLOCO 7 · O FECHO                                       slides 40–44
# ══════════════════════════════════════════════════════════════════════

# ── 42 · o ciclo completo ─────────────────────────────────────────────
s, y = cabecalho("O que existe agora", "O ciclo fechado")
h = mono(s, y + 0.2,
     "  CSV  →  bronze  →  silver  →  gold  →  modelo  →  fila  →  app\n"
     "                                  ↑                              │\n"
     "                                  │        retorno da ligação    │\n"
     "                                  └──────────────────────────────┘", tam=14)
bloco(s, y + h + 0.45,
      "Quatro noites. Um catálogo, um pipeline de 16 tarefas, um modelo, dois "
      "Genies e um app.\nE o dado dá a volta inteira.",
      h=1.2, tam=16.5, alinha=PP_ALIGN.CENTER)

# ── 43 · as quatro noites em quatro números ───────────────────────────
s, y = cabecalho("Uma query só", "A imersão inteira, em quatro linhas")
tabela(s, y + 0.1,
       ["ETAPA", "NÚMERO", "O QUE É"],
       [["noite 1 · o dado", "28.729", "pedidos que chegaram do ERP"],
        ["noite 2 · o pipeline", "191.080", "linhas em gold.fato_vendas"],
        ["noite 3 · a decisão", "200", "contatos na fila desta semana"],
        ["noite 4 · o retorno", "1", "ligações que o time registrou"]],
       larguras=[4.2, 2.4, 5.233], alt=0.62, destaque=3)
bloco(s, 6.0, "A última linha só tem número porque alguém clicou num botão "
               "aqui, agora.", h=0.7, tam=16)

# ── 44 · o que você leva ──────────────────────────────────────────────
s, y = cabecalho("Não é um curso que você assistiu", "O que fica com você")
cartoes(s, y + 0.1, 2.5, [
    ("Um repositório",
     "Bundle versionado, 16 tarefas, testes que quebram o job, dois Genies e um "
     "app. Tudo em código, tudo no Git."),
    ("Um caso de negócio",
     "“4,25× mais pedidos com a mesma equipe” é uma frase de entrevista. E você "
     "sabe explicar cada número dela."),
    ("Um método",
     "Prompt, deploy, conferência contra o banco. Repetido quinze vezes em "
     "quatro noites.", "destaque"),
])
bloco(s, 5.9, "Rode o 99-limpar de cada noite e refaça sozinho. É aí que "
               "assenta.", h=0.7, tam=16)

# ── 45 · o que estudar depois ─────────────────────────────────────────
s, y = cabecalho("Se você quer continuar", "Por onde seguir")
faixas(s, y + 0.05, [
    ("Refazer as quatro noites sozinho", "os scripts 99-limpar devolvem o ambiente ao estado anterior"),
    ("Trocar o dataset pelo seu", "o formato é CSV; a sujeira do seu ERP é mais interessante que a nossa"),
    ("Agendar o job e monitorar", "o pipeline já é agendado — falta você olhar o que ele faz sem você"),
    ("Retreinar com o retorno real", "daqui a um mês, gold.retorno_ligacao tem dado suficiente para virar rótulo", "destaque"),
], alt=0.66, larg_esq=5.6)
bloco(s, 6.2, "O último item é o projeto ficando mais inteligente sozinho — "
               "porque o ciclo fechou.", h=0.66, tam=15.5)

# ── 46 · a documentação, para depois da aula ──────────────────────────
s, y = cabecalho("Para ler com calma amanhã", "A documentação do que a gente usou",
                 "Tudo desta noite está documentado. Estes são os pontos de "
                 "partida — o resto se acha a partir deles.")
faixas(s, y + 0.05, [
    ("Databricks Apps", "docs.databricks.com/dev-tools/databricks-apps · o que é, frameworks e deploy"),
    ("Apps · conceitos e permissões", ".../databricks-apps/key-concepts · o service principal e o modelo de acesso", "destaque"),
    ("Apps · boas práticas", ".../databricks-apps/best-practices"),
    ("AI/BI Genie", "docs.databricks.com/genie · a experiência de perguntar em português"),
    ("Genie Agents · afinar", "docs.databricks.com/genie-agents/tune-quality · instruções e qualidade da resposta"),
    ("Unity Catalog · privilégios", "docs.databricks.com/data-governance/unity-catalog/manage-privileges · o GRANT"),
], alt=0.6, larg_esq=4.6, cabecalho_=("O ASSUNTO", "ONDE LER"))
nota(s, "Links conferidos contra o índice oficial (docs.databricks.com/llms.txt) "
        "no dia da aula.")

# ── 47 · a frase da noite ─────────────────────────────────────────────
impacto("Segunda, uma query quebrou.\n"
        "Hoje, um clique virou uma linha na gold.\n\n"
        "E vocês construíram junto. Não assistiram.")

# ══════════════════════════════════════════════════════════════════════
import os
saida = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                     "aula-04-app-e-genie.pptx")
prs.save(saida)
print(f"{len(prs.slides._sldIdLst)} slides → {saida}")
